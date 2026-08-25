#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build_pptx_kr.py — 한국형 PPTX IM (D19 §6 개정 · KR 판)

실제 국내 IM 3종(당산동·잠원동·Basic)을 분해해 관행을 반영했습니다.

  · **사진 전진배치** — 표지는 외관 풀블리드, 개요는 좌 스펙표 + 우 대형 사진
  · **개조식 카피** — 슬라이드 불릿은 명사형. 완결 문장은 주석·경고에만
  · **매매가는 적색 강조** — 국내 IM의 사실상 표준
  · 라벨 열 남색 바탕 흰 글씨 · 값 우측 정렬 · 강조값 적색
  · 지도·지적도 면 · 사진 갤러리 2면

정직성 규칙은 그대로입니다 — 확인 필요·기준 병기·역레버리지 경고·출처 표기.
"""
from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

import core as im_core
import public_data as PD
import input_spec as ISPEC
from parcel import PROV_KO
import presets as PS
from core import ASSUMED_LOAN_RATE, MAN, PYEONG, eok, man, pct

sys.path.insert(0, str(Path(__file__).resolve().parent / 'ssot'))
import loader as SSOT                                            # noqa: E402

SSOT_PARCEL = SSOT.load('im.parcel')

HERE = Path(__file__).resolve().parent
SW, SH = 13.333, 7.5
M, CW = 0.60, 12.133
RIGHT, BODY_Y, FOOT_Y = M + CW, 6.86, 7.00
FONT = 'Noto Sans CJK KR'


def C(h):
    return RGBColor(*(int(h[i:i + 2], 16) for i in (0, 2, 4)))


K: dict = {}
PRESET: PS.Preset = PS.P6


def apply_preset(p: PS.Preset) -> None:
    """아키타입은 그대로, 토큰만 갈아 끼웁니다 (D19 §6B)."""
    global K, PRESET
    PRESET = p
    t = p.tokens
    K.clear()
    K.update({
        'navy':   C(t['labelCol'].lstrip('#')),
        'navy_d': C(t['accentD'].lstrip('#')),
        'ink':    C(t['body'].lstrip('#')),
        'ink2':   C(t['ink3'].lstrip('#')),
        'mute':   C(t['mute'].lstrip('#')),
        'line':   C(t['line'].lstrip('#')),
        'panel':  C(t['tint'].lstrip('#')),
        'white':  C('FFFFFF'),
        'red':    C(t['priceHi'].lstrip('#')),
        'redbg':  C(t['priceHiBg'].lstrip('#')),
        'green':  C(t['green'].lstrip('#')),
        'greenbg': C(t['greenL'].lstrip('#')),
        'gold':   C(t['accent'].lstrip('#')),
        'chip':   C(t['srcChip'].lstrip('#')),
        'chipbg': C(t['srcChipBg'].lstrip('#')),
    })


apply_preset(PS.P6)


class Overflow(Exception):
    pass


# ── 원시 도형 ──────────────────────────────────────────────────────────
def rect(s, x, y, w, h, fill=None, line=None, lw=0.75):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                            Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    sp.text_frame.text = ''
    return sp


def txt(s, x, y, w, h, text, *, size=11.5, bold=False, color='ink',
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=0, wrap=True,
        line_spacing=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = K[color] if isinstance(color, str) else color
    return tb


def pick(A: dict, *keys):
    """있는 사진을 순서대로 찾습니다. 없으면 None."""
    for k in keys:
        if k in A:
            return A[k]
    return None


def photo_or_note(s, A, keys, x, y, w, h, note='사진 미제출'):
    src = pick(A, *keys)
    if src is None:
        rect(s, x, y, w, h, fill=K['panel'], line=K['line'])
        txt(s, x, y + h / 2 - 0.16, w, 0.30, note, size=11.5, color='mute',
            align=PP_ALIGN.CENTER)
        return None
    return photo(s, src, x, y, w, h)


def photo(s, path: Path, x, y, w, h):
    """지정한 상자를 꽉 채우도록 잘라 넣습니다 (object-fit: cover)."""
    im = Image.open(path)
    box_ar, img_ar = w / h, im.width / im.height
    if img_ar > box_ar:                       # 가로가 길다 → 좌우를 자른다
        nw = int(im.height * box_ar)
        im = im.crop(((im.width - nw) // 2, 0, (im.width + nw) // 2, im.height))
    else:                                     # 세로가 길다 → 위아래를 자른다
        nh = int(im.width / box_ar)
        top = int((im.height - nh) * 0.30)    # 건물 몸통이 중앙에 오도록
        im = im.crop((0, top, im.width, top + nh))
    tmp = Path('/tmp/_pp.jpg')
    im.convert('RGB').save(tmp, 'JPEG', quality=88)
    return s.shapes.add_picture(str(tmp), Inches(x), Inches(y),
                                Inches(w), Inches(h))


def chip(s, x, y, text, *, fill='navy', color='white', size=10.5, pad=0.13):
    w = len(text) * size / 72 * 1.06 + pad * 2
    rect(s, x, y, w, 0.28, fill=K[fill])
    txt(s, x + pad, y + 0.045, w - pad * 2, 0.20, text, size=size,
        bold=True, color=color)
    return w


def table(s, x, y, w, head, rows, widths, *, size=10, head_size=10,
          row_h=0.30, head_h=0.32, align=None, label_col=False,
          accent_rows=None, head_fill='navy'):
    """국내 IM 관행 — 머리행 남색 바탕 흰 글씨, 강조 행은 적색."""
    h = head_h + row_h * len(rows)
    if y + h > BODY_Y + 1e-6:
        raise Overflow(f'표가 본문 영역을 넘습니다: {y + h:.2f} > {BODY_Y}')
    shp = s.shapes.add_table(len(rows) + 1, len(head), Inches(x), Inches(y),
                             Inches(w), Inches(h))
    t = shp.table
    t.first_row = True
    t.horz_banding = False
    tot = sum(widths)
    for i, cw in enumerate(widths):
        t.columns[i].width = Emu(int(Inches(w) * cw / tot))
    t.rows[0].height = Inches(head_h)
    for i in range(len(rows)):
        t.rows[i + 1].height = Inches(row_h)
    AL = {'l': PP_ALIGN.LEFT, 'r': PP_ALIGN.RIGHT, 'c': PP_ALIGN.CENTER}
    align = align or ['l'] * len(head)
    accent_rows = accent_rows or {}

    def put(cell, s_, *, bold, sz, col, al, bg):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        cell.margin_left = cell.margin_right = Inches(0.07)
        cell.margin_top = cell.margin_bottom = Inches(0.01)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = AL[al]
        r = p.add_run()
        r.text = str(s_)
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = col

    for j, hd in enumerate(head):
        put(t.cell(0, j), hd, bold=True, sz=head_size,
            col=K['white'] if head_fill == 'navy' else K['ink2'],
            al='c', bg=K[head_fill])
    for i, row in enumerate(rows):
        acc = accent_rows.get(i)
        for j, v in enumerate(row):
            sv = str(v)
            neg = sv.startswith('▼') or sv.startswith('−') or sv.startswith('-')
            col = K['red'] if (acc == 'red' or neg) else K['ink']
            bold = bool(acc) or neg or (label_col and j == 0)
            bg = (K['redbg'] if acc == 'red'
                  else K['navy'] if (label_col and j == 0)
                  else K['panel'] if acc == 'sum' else K['white'])
            if label_col and j == 0:            # 라벨 열은 강조 행에서도 남색
                col, bold, bg = K['white'], True, K['navy']
            put(t.cell(i + 1, j), sv, bold=bold, sz=size, col=col,
                al=align[j], bg=bg)
    return shp


# ── 지면 틀 ────────────────────────────────────────────────────────────
def page(prs, kicker_en, title_ko, *, badge=None, badge_kind='ok'):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    txt(s, M, 0.30, CW, 0.20, kicker_en, size=9, color='mute')
    rect(s, M, 0.585, 0.20, 0.30, fill=K['navy'])
    txt(s, M + 0.30, 0.55, CW - 3.4, 0.38, title_ko, size=21, bold=True,
        color='navy')
    if badge:
        bg, fg = (('greenbg', 'green') if badge_kind == 'ok' else ('redbg', 'red'))
        w = len(badge) * 0.105 + 0.34
        rect(s, RIGHT - w, 0.58, w, 0.32, fill=K[bg])
        txt(s, RIGHT - w, 0.645, w, 0.22, badge, size=10.5, bold=True,
            color=fg, align=PP_ALIGN.CENTER)
    rect(s, M, 1.00, CW, 0.022, fill=K['navy'])
    s._chip_band = False
    return s


def chip_band(s, core, y=1.06):
    """evidence 레이아웃 — 매 면 헤더에 근거 칩 (P7 §switches.sourceChips)."""
    if PRESET.switches.get('sourceChips') != 'header':
        return 0.0
    x = M
    for c in core.source_chips()[:6]:
        w = len(c) * 0.082 + 0.22
        rect(s, x, y, w, 0.24, fill=K['chipbg'])
        txt(s, x, y + 0.045, w, 0.18, c, size=8.5, bold=True,
            color=K['chip'], align=PP_ALIGN.CENTER)
        x += w + 0.08
    return 0.34


def _wrap45(t: str, n: int = 44) -> list[str]:
    """45자를 넘지 않게 나눕니다. 낱말 경계를 지킵니다 (정본 §4.7)."""
    out, cur = [], ''
    for w in t.split(' '):
        if len(cur) + len(w) + 1 > n:
            out.append(cur.strip())
            cur = w
        else:
            cur += (' ' if cur else '') + w
    if cur.strip():
        out.append(cur.strip())
    return out


def chips(s, x0: float, y: float, items: list[str], size: float = 9,
          maxw: float = 7.10) -> float:
    """작은 라벨 칩을 가로로 늘어놓습니다. 넘치면 다음 줄로 접습니다."""
    x, y0 = x0, y
    for c in items:
        w = len(c) * (size * 0.0105) + 0.20
        if x + w > x0 + maxw:
            x, y = x0, y + 0.30
        rect(s, x, y, w, 0.26, fill=K['chipbg'])
        txt(s, x, y + 0.05, w, 0.20, c, size=size, bold=True,
            color=K['chip'], align=PP_ALIGN.CENTER)
        x += w + 0.08
    return y - y0 + 0.30


def footer(prs, core, org='제이에스부동산중개(주)'):
    n = len(prs.slides)
    for i, s in enumerate(prs.slides, 1):
        if i == 1:
            continue
        txt(s, M, FOOT_Y, CW * 0.7, 0.22,
            f'{org}   |   본 자료는 투자 권유가 아니며 실사·전문가 검토가 필요합니다',
            size=8, color='mute', wrap=False)
        txt(s, M + CW * 0.7, FOOT_Y, CW * 0.3, 0.22, f'{i} / {n}',
            size=8, color='mute', align=PP_ALIGN.RIGHT)


# ── 슬라이드 ───────────────────────────────────────────────────────────
def s01_cover(prs, core, A):
    """표지 — 좌 정보 패널 + 우 외관 사진 전면.

    풀블리드로 깔면 4:3 원본이 위아래로 잘려 건물 아랫부분이 사라집니다.
    좌우 분할이면 세로를 살린 채 건물 전체가 보입니다.
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    style = PRESET.cover
    PW = 7.60 if style != 'land_metric' else 8.20
    px = SW - PW
    photo_or_note(s, A, ['hero', 'front'], px, 0, PW, SH)
    rect(s, 0, 0, px, SH, fill=K['navy_d'])
    rect(s, px - 0.07, 0, 0.07, SH, fill=K['gold'])

    gu = core.address_band.split()[-1]
    alias = COPY[core.fixture_id]['alias']
    L = 0.72
    W = px - L - 0.5
    txt(s, L, 1.32, W, 0.24, '· INFORMATION MEMORANDUM', size=10.5,
        color=RGBColor(0xC9, 0xB3, 0x86))
    rect(s, L, 1.72, 0.62, 0.30, fill=K['gold'])
    txt(s, L, 1.775, 0.62, 0.22, '매각', size=11, bold=True, color='navy_d'
        if False else 'white', align=PP_ALIGN.CENTER)
    txt(s, L, 2.24, W, 1.32, f'{gu}\n{alias} 매각', size=34, bold=True,
        color='white', space=6, line_spacing=1.15)
    txt(s, L, 3.66, W, 0.26, core.address_band, size=12.5,
        color=RGBColor(0xA9, 0xB6, 0xC8))
    if core.has_public:
        st = core.nearest_station
        txt(s, L, 3.96, W, 0.26,
            f'{st[0].split(" (")[0]} {st[1]}m  ·  {core.f("zoning").value}',
            size=12.5, color=RGBColor(0xA9, 0xB6, 0xC8))

    if style == 'land_metric' and core.has_public:
        items = [('매매 희망가', eok(core.price, 0) + '원', True),
                 ('토지 평당가',
                  f'{core.land_pyeong_price.value / MAN:,.0f}만원/평', False),
                 ('공시지가 배수', f'{core.land_price_multiple.value:.2f}배', False),
                 ('잔여 용적률', f'{core.far_headroom.value:.1f}%p', False)]
    else:
        items = [('매매 희망가', eok(core.price, 0) + '원', True),
                 ('토지 평당가',
                  f'{core.land_pyeong_price.value / MAN:,.0f}만원/평'
                  if core.land_pyeong_price.known else '확인 필요', False),
                 ('월 임대료', man(core.monthly_rent) + '원', False)]
    if style == 'evidence_band':
        cx = L
        for c in core.source_chips()[:4]:
            cw2 = len(c) * 0.085 + 0.24
            rect(s, cx, 5.98, cw2, 0.26, fill=K['chipbg'])
            txt(s, cx, 6.03, cw2, 0.20, c, size=8.5, bold=True,
                color=K['chip'], align=PP_ALIGN.CENTER)
            cx += cw2 + 0.09
    y = 4.30 if len(items) > 3 else 4.48
    step = 0.54 if len(items) > 3 else 0.62
    for lb, v, hot in items:
        rect(s, L, y, W, 0.03, fill=RGBColor(0x2A, 0x3D, 0x5C))
        txt(s, L, y + 0.14, 2.0, 0.24, lb, size=10.5,
            color=RGBColor(0x9F, 0xAF, 0xC5))
        txt(s, L + 1.95, y + 0.06, W - 1.95, 0.40, v, size=20, bold=True,
            color=K['red'] if hot else K['white'], align=PP_ALIGN.RIGHT)
        y += step

    txt(s, L, 6.52, W, 0.26, '제이에스부동산중개(주)', size=12.5, bold=True,
        color='white')
    txt(s, L, 6.82, W, 0.22, '2026. 08   ·   본 자료는 투자 권유가 아닙니다',
        size=9.5, color=RGBColor(0x8B, 0x9B, 0xB2))
    return s


def fill_tokens(items: list[str], core) -> list[str]:
    """손으로 쓴 카피 안의 수치 자리표시자를 코어 값으로 채웁니다.

    🔴 카피에 숫자를 손으로 적으면 반드시 어긋납니다. 실제로 이 덱의
       요약 면에 "1억 5,923만원" 이 적혀 있었고 표지·개요·토지가치 면은
       "15,933만원" 이었습니다. 같은 IM 안에서 평당가가 두 값이었습니다.
    """
    lp = core.land_pyeong_price
    tok = {'land_pyeong': (f'{lp.value / MAN:,.0f}만원' if lp.known else '확인 필요')}
    return [it.format(**tok) if '{' in it else it for it in items]


def s02_points(prs, core, A, pts):
    pts = fill_tokens(pts, core)
    """제안 Point — 좌 사진 2장, 우 개조식 불릿."""
    s = page(prs, 'Executive Summary', '한 장 요약 · 제안 Point')
    photo_or_note(s, A, ['front','hero'], M, 1.24, 5.10, 2.62)
    photo_or_note(s, A, ['ext1','context'], M, 3.98, 2.48, 1.92)
    photo_or_note(s, A, ['ext2','map'], M + 2.62, 3.98, 2.48, 1.92)

    x = M + 5.42
    w = RIGHT - x
    for i, p in enumerate(pts):
        y = 1.30 + i * 0.60
        rect(s, x, y + 0.055, 0.055, 0.28, fill=K['navy'])
        txt(s, x + 0.20, y, w - 0.20, 0.52, p, size=13, color='ink2',
            space=0, line_spacing=1.25)

    y = 1.30 + len(pts) * 0.60 + 0.10
    rect(s, x, y, w, 1.02, fill=K['panel'])
    ltv = {int(r['ltv'] * 100): r for r in core.ltv_rows}[50]
    txt(s, x + 0.22, y + 0.16, w - 0.44, 0.24,
        '· 실투자금 · 월 순현금 (LTV 50% · 금리 4.5% 가정)', size=10.5,
        color='mute')
    txt(s, x + 0.22, y + 0.46, (w - 0.44) / 2, 0.40,
        eok(ltv['equity']), size=19, bold=True)
    sign = '▼ ' if ltv['monthly_net'] < 0 else '+'
    txt(s, x + 0.22 + (w - 0.44) / 2, y + 0.46, (w - 0.44) / 2, 0.40,
        f'{sign}{man(ltv["monthly_net"])}원/월', size=19, bold=True,
        color='red' if ltv['monthly_net'] < 0 else 'ink')

    y2 = y + 1.16
    rect(s, x, y2, w, 1.42, fill=None, line=K['line'])
    rect(s, x, y2, 0.06, 1.42, fill=K['red'])
    txt(s, x + 0.24, y2 + 0.14, w - 0.48, 0.24,
        f'확인 필요 {len(core.deficiencies)}건', size=11.5, bold=True,
        color='red')
    for i, d in enumerate(core.deficiencies[:4]):
        txt(s, x + 0.24, y2 + 0.44 + i * 0.24, w - 0.48, 0.22,
            '· ' + d.split(' (')[0].split(' —')[0], size=10, color='ink2')
    return s


def s03_overview(prs, core, A, tags, bullets):
    """물건 개요 — 좌 스펙표 + 우 대형 사진. 국내 IM의 기본형."""
    s = page(prs, 'Property Overview', '물건 개요',
             badge='공부 확인' if core.has_public else '중개인 제공',
             badge_kind='ok' if core.has_public else 'warn')
    chip_band(s, core, 1.06)
    tw = 5.55
    rows, acc = [], {}
    land, gfa, pp = core.land_sqm, core.gfa_confirmed, core.land_pyeong_price
    P_ = lambda v: f'{v:,.2f}㎡ ({v / PYEONG:,.1f}평)'          # noqa: E731
    rows.append(['소재지', core.address_band])
    if core.has_public:
        f = core.f
        rows += [['용도지역', f('zoning').value],
                 ['대지면적', P_(land.value)],
                 ['건축면적', P_(f('archSqm').value)],
                 ['연면적', P_(gfa.value)],
                 ['건폐율 / 용적률',
                  f'{f("bcrPct").value:.2f}% / {f("farPct").value:.2f}%'
                  f'  (상한 {f("bcrLimit").value:.0f}% / {f("farLimit").value:.0f}%)'],
                 ['주용도 / 구조', f'{f("mainUse").value} · {f("structure").value}'],
                 ['규모', f('floors').value],
                 ['사용승인', f'{f("approvalDate").value}  ·  위반건축물 {f("violation").value}'],
                 ['주차 / E·V', f'{f("parking").value} · {f("elevator").value}'],
                 ['공시지가',
                  f'{f("landPriceSqm").value * PYEONG / MAN:,.0f}만원/평  '
                  f'(총 {eok(core.land_price_total.value)})']]
    else:
        rows += [['용도지역', '확인 필요'], ['대지면적', '확인 필요'],
                 ['연면적', '확인 필요 — 계 행과 층별 합 불일치'],
                 ['건폐율 / 용적률', '확인 필요'],
                 ['주용도', core.building_use],
                 ['규모', f'{len(core.rows)}개 구획'],
                 ['공시지가', '확인 필요']]
    rows.append(['매매 희망가', eok(core.price, 0) + '원'])
    acc[len(rows) - 1] = 'red'
    rows.append(['토지 평당가',
                 f'{pp.value / MAN:,.0f}만원/평' if pp.known else '확인 필요'])
    acc[len(rows) - 1] = 'sum'
    table(s, M, 1.22, tw, ['구  분', '내       용'], rows, [1.55, 4.0],
          size=10, row_h=0.335, align=['c', 'l'], label_col=True,
          accent_rows=acc)

    px = M + tw + 0.22
    pw = RIGHT - px
    photo_or_note(s, A, ['front', 'ext1', 'hero'], px, 1.22, pw, 3.44)
    cx = px + 0.14
    for t in tags:
        cx += chip(s, cx, 1.36, t) + 0.10
    for i, b in enumerate(bullets):
        y = 4.84 + i * 0.40
        rect(s, px, y + 0.09, 0.05, 0.16, fill=K['navy'])
        txt(s, px + 0.18, y, pw - 0.18, 0.38, b, size=11.5, color='ink2')

    # 매각 조건 띠 — 국내 IM 관행
    y0 = 6.10
    rect(s, M, y0, CW, 0.62, fill=K['panel'])
    rect(s, M, y0, 1.30, 0.62, fill=K['navy'])
    txt(s, M, y0 + 0.18, 1.30, 0.26, '매각 조건', size=11.5, bold=True,
        color='white', align=PP_ALIGN.CENTER)
    if core.has_public and core.gfa_confirmed.note:
        txt(s, M, 5.72, tw, 0.32,
            '· 연면적은 건축물대장 층별개요 합으로 확정 (용적률 역산 검증)\n'
            f'· 원장 계 행 {core.stated_area:,.2f}㎡는 오기 — 확인 필요',
            size=9.5, color='mute', space=1)
    cells = [('매매 희망가', eok(core.price, 0) + '원', 'red'),
             ('승계 보증금', eok(core.deposit) + '원', 'ink'),
             ('월 임대료', man(core.monthly_rent) + '원', 'ink'),
             ('연 수익률', f'{pct(core.gross_price.value)} (총임대료 ÷ 매매가)', 'ink')]
    cw2 = (CW - 1.30) / 4
    for i, (lb, v, col) in enumerate(cells):
        x2 = M + 1.30 + i * cw2
        txt(s, x2 + 0.18, y0 + 0.10, cw2 - 0.3, 0.20, lb, size=9.5, color='mute')
        txt(s, x2 + 0.18, y0 + 0.31, cw2 - 0.3, 0.26, v, size=13, bold=True,
            color=col)
    return s


def s04_location(prs, core, A, notes):
    s = page(prs, 'Location', '입지 · 교통',
             badge='카카오 로컬 조회' if core.has_public else '확인 필요',
             badge_kind='ok' if core.has_public else 'warn')
    photo_or_note(s, A, ['map'], M, 1.22, 8.05, 5.50, '지도 미확보')
    x = M + 8.25
    w = RIGHT - x
    if core.has_public:
        tr = core.f('transit').value
        rows = [[n.split(' (')[0], f'{d:,}m', t] for n, d, t in tr]
        table(s, x, 1.22, w, ['구분', '거리', '도보'], rows, [2.0, 1.0, 1.1],
              size=10, row_h=0.33, align=['l', 'r', 'c'])
        y = 1.22 + 0.32 + 0.33 * len(rows) + 0.28
    else:
        rect(s, x, 1.22, w, 1.0, fill=K['redbg'])
        txt(s, x + 0.16, 1.42, w - 0.32, 0.62,
            '인근 역까지의 실측 거리가\n제출 자료에 없습니다', size=11.5,
            color='red', space=2)
        y = 2.42
    for i, n in enumerate(notes):
        rect(s, x, y + i * 0.62, 0.05, 0.20, fill=K['navy'])
        txt(s, x + 0.18, y + i * 0.62 - 0.03, w - 0.18, 0.58, n, size=11,
            color='ink2', line_spacing=1.2)
    txt(s, M, 6.80, CW, 0.20,
        '지도 출처 — 카카오맵 (권역 참고용) · 정확한 위치는 현장 안내 · 거리는 위 표 기준',
        size=8.5, color='mute')
    return s


def s05_land(prs, core, A):
    s = page(prs, 'Land & Zoning', '토지 이용 · 공법 규제',
             badge='토지이용계획 조회' if core.has_public else '확인 필요',
             badge_kind='ok' if core.has_public else 'warn')
    zv = core.zoning_view()
    ph_h = 5.10 if not zv['본문'] else 3.70
    photo_or_note(s, A, ['cadastre','map'], M, 1.22, 7.10, ph_h, '지적도 미확보')
    if zv['본문']:
        # L12 — 매수 목적에 맞는 항목만 본문에. 전체는 부록에 실립니다.
        yz = 1.22 + ph_h + 0.20
        lbl = (SSOT_PARCEL['zoning_display']['purpose_labels']
               .get(core.buyer_purpose, '목적 미지정'))
        txt(s, M, yz, 7.10, 0.26,
            f'· 토지이용계획 · {lbl} 관점 {len(zv["본문"])}항목', size=11,
            color='mute')
        txt(s, M, yz + 0.28, 7.10, 0.26,
            f'· 전체 {len(zv["부록"])}항목은 부록', size=11, color='mute')
        chips(s, M, yz + 0.58, zv['본문'][:6], size=11)
    x = M + 7.30
    w = RIGHT - x
    if core.has_public:
        f = core.f
        rows = [['용도지역', f('zoning').value],
                ['기타 지구', f('zoningOverlap').value.replace(' · ', '\n')],
                ['법정 건폐율', f'{f("bcrLimit").value:.0f}%  (현행 {f("bcrPct").value:.2f}%)'],
                ['법정 용적률', f'{f("farLimit").value:.0f}%  (현행 {f("farPct").value:.2f}%)']]
        table(s, x, 1.22, w, ['구  분', '내  용'], rows, [1.4, 3.0],
              size=9.5, row_h=0.42, align=['c', 'l'], label_col=True)
        hr = core.far_headroom
        y = 1.22 + 0.32 + 0.42 * 4 + 0.30
        fill_bg = 'greenbg' if hr.value > 20 else 'panel'
        rect(s, x, y, w, 1.30, fill=K[fill_bg])
        txt(s, x + 0.18, y + 0.16, w - 0.36, 0.22, '잔여 용적률', size=10.5,
            color='mute')
        txt(s, x + 0.18, y + 0.42, w - 0.36, 0.48,
            f'{hr.value:.1f}%p', size=26, bold=True,
            color='green' if hr.value > 20 else 'ink')
        txt(s, x + 0.18, y + 0.96, w - 0.36, 0.24,
            '· 증축 규모는 높이·주차 기준 검토 후 산출' if hr.value > 20
            else '· 법정 상한까지 개발 완료', size=10, color='ink2')

        # ── L11 · 제척이 있으면 유효 면적을 병기합니다 ──
        # 🔴 대장 면적만 쓰면 증축 여유를 과대평가합니다.
        L = core.land
        if L and L.excluded_area > 0:
            y2 = y + 1.44
            rect(s, x, y2, w, 1.62, fill=K['redbg'])
            txt(s, x + 0.18, y2 + 0.14, w - 0.36, 0.22,
                f'제척 {L.excluded_area:,.1f}㎡ 반영 — 유효 기준', size=11.5,
                bold=True, color='red')
            rows2 = [['대장 대지', f'{L.ledger_area:,.1f}㎡',
                      f'{L.ledger_far_pct:.1f}%'],
                     ['유효 대지', f'{L.effective_area:,.1f}㎡',
                      f'{L.effective_far_pct:.1f}%']]
            table(s, x + 0.14, y2 + 0.40, w - 0.28,
                  ['구  분', '대지면적', '용적률'], rows2, [1.0, 1.1, 1.0],
                  size=9, row_h=0.32, align=['c', 'r', 'r'], label_col=True)
            rc = core.relief_cross
            txt(s, x + 0.18, y2 + 1.30, w - 0.36, 0.28,
                (f'▲ 한시 완화 {rc["threshold_pct"]:.0f}% 조건 이탈'
                 if rc else '· 제척은 ●중개인 · 구청 확인 필요'),
                size=11, bold=bool(rc), color='red' if rc else 'mute')
    else:
        rect(s, x, 1.22, w, 1.2, fill=K['redbg'])
        txt(s, x + 0.16, 1.44, w - 0.32, 0.80,
            '용도지역·건폐율·용적률\n확인 필요\n(토지이용계획 미조회)',
            size=12, bold=True, color='red', space=2)
    return s


def s06_rentroll(prs, core, A):
    s = page(prs, 'Rent Roll', '층별 이용 및 임대 현황',
             badge=f'확인 필요 {len(core.deficiencies)}건', badge_kind='warn')
    rows, acc = [], {}
    for r in core.rows:
        a = r.get('leaseAreaSqm')
        grp = f' ({r["contractGroup"]})' if r.get('contractGroup') else ''
        rows.append([
            r['unitLabel'] + grp,
            f'{a:,.2f}' if a else '확인 필요',
            f'{a / PYEONG:,.1f}' if a else '—',
            r['tenantBusiness'] or '자가 / 공실',
            f"{r['depositKrw'] / MAN:,.0f}" if r['depositKrw'] else '—',
            f"{r['monthlyRentKrw'] / MAN:,.0f}" if r['monthlyRentKrw'] else '—',
            r['currentExpiryDate'] or '확인 필요',
            r['leaseState'],
        ])
    tot_a = core.ledger_sum_area
    rows.append(['계', f'{tot_a:,.2f}' if tot_a else '—',
                 f'{tot_a / PYEONG:,.1f}' if tot_a else '—', '',
                 f'{core.ledger_sum_deposit / MAN:,.0f}',
                 f'{core.ledger_sum_rent / MAN:,.0f}', '', ''])
    acc[len(rows) - 1] = 'sum'
    table(s, M, 1.22, CW,
          ['층별', '임대면적(㎡)', '(평)', '임차 업종', '보증금(만원)',
           '월세(만원)', '계약 만료일', '상태'],
          rows, [1.0, 1.5, 1.0, 2.4, 1.6, 1.5, 1.7, 1.1],
          size=9.5, row_h=0.295, align=['c', 'r', 'r', 'l', 'r', 'r', 'c', 'c'],
          accent_rows=acc)
    y = 1.22 + 0.32 + 0.295 * len(rows) + 0.24
    st = core.state_counts
    bu, ba = core.vacancy
    f1 = core.first_floor_share
    notes = [
        '구성 — ' + ' · '.join(f'{k} {v}' for k, v in st.items()),
        f'공실률 — 구획 {pct(bu.value)} · 면적 '
        f'{pct(ba.value) if ba.known else "확인 필요"}',
        (f'1층 월세 비중 {pct(f1.value, 1)} — 소형상업용은 1층이 수입의 절반 안팎을 좌우'
         if f1.known else '1층 비중 확인 필요'),
        '금액은 VAT 별도 · 임차인 상호는 대외 문서에 표기하지 않음',
    ]
    for i, n in enumerate(notes):
        txt(s, M, y + i * 0.28, CW, 0.26, '· ' + n, size=10.5, color='ink2')

    y2 = y + len(notes) * 0.28 + 0.22
    if y2 + 1.05 <= BODY_Y:
        rect(s, M, y2, CW, 1.02, fill=K['redbg'])
        rect(s, M, y2, 0.06, 1.02, fill=K['red'])
        txt(s, M + 0.24, y2 + 0.14, CW - 0.48, 0.26,
            '매입 후 수익 개선 여지', size=12, bold=True, color='red')
        txt(s, M + 0.24, y2 + 0.44, CW - 0.48, 0.50,
            COPY[core.fixture_id]['upside_line'], size=11, color='ink2',
            space=2, line_spacing=1.2)
    return s


def s07_lease2(prs, core, A, upside):
    s = page(prs, 'Lease Analysis', '임대차 분석 · 개선 여력')
    left = 6.05
    rows = [['갱신요구권', '산출하지 않음 — 최초 계약일 미기재'],
            ['환산보증금', '전 계약 서울 기준 9억 이하 · 상가임대차법 전면 적용'],
            ['만료 임박', ''], ['운영비', '자료 미제출 — 순수입 기준 수익률 미산출']]
    ref = core.as_of or '2026-08-24'
    from datetime import date as _d
    r0 = _d.fromisoformat(ref)
    exp = [r for r in core.rows if core.expiry_state(r, r0) == '만료 경과']
    soon = [r for r in core.rows if core.expiry_state(r, r0) == '만료 임박']
    rows[2][1] = (f'만료 경과 {len(exp)}건 · 60일 내 {len(soon)}건'
                  if exp or soon else '해당 없음')
    table(s, M, 1.22, left, ['항  목', '판  정'], rows, [1.5, 4.0],
          size=10, row_h=0.44, align=['c', 'l'], label_col=True)

    x = M + left + 0.28
    w = RIGHT - x
    txt(s, x, 1.22, w, 0.28, '개선 여력', size=13, bold=True, color='navy')
    for i, u in enumerate(upside):
        y = 1.60 + i * 0.62
        rect(s, x, y + 0.08, 0.05, 0.20, fill=K['gold'])
        txt(s, x + 0.18, y, w - 0.18, 0.58, u, size=11.5, color='ink2',
            line_spacing=1.2)

    y = 1.22 + 0.32 + 0.44 * 4 + 0.34
    rect(s, M, y, left, 1.20, fill=K['panel'])
    txt(s, M + 0.18, y + 0.16, left - 0.36, 0.24,
        '· 자료 해상도 — 임대차 L축 / 물건자료 P축', size=10.5, color='mute')
    Lg, Pg, short2 = core.resolution_pair()
    txt(s, M + 0.18, y + 0.42, 1.60, 0.40, f'{Lg} · {Pg}', size=22, bold=True,
        color='navy')
    nxt = next((g for g in ('R1', 'R2', 'R3', 'P1', 'P2', 'P3')
                if short2.get(g)), None)
    miss = short2.get(nxt) or []
    txt(s, M + 1.90, y + 0.44, left - 2.10, 0.60,
        ('· 다음 등급 ' + nxt + ' 까지 부족\n· ' +
         ' · '.join(ISPEC.label(k) for k in miss)) if miss else '· 전 항목 확보',
        size=10.5, color='ink2', line_spacing=1.15)
    return s


def s08_invest(prs, core, A):
    s = page(prs, 'Investment Structure', '투자 구조 — 실투자금과 월 순현금')
    left = 5.35
    tac = core.total_acq_cost.value
    rows = [['매매가', eok(core.price)],
            ['취득세 4.6%', eok(core.acq_tax.value)],
            ['중개보수 0.9%', eok(core.broker_fee.value)],
            ['기타비용', '자료 없음'],
            ['총취득원가', eok(tac)]]
    table(s, M, 1.22, left, ['구  분', '금  액'], rows, [2.0, 2.2],
          size=10.5, row_h=0.36, align=['c', 'r'], label_col=True,
          accent_rows={4: 'red'})
    txt(s, M, 1.22 + 0.32 + 0.36 * 5 + 0.14, left, 0.24,
        f'· 매매가보다 {eok(tac - core.price)}원({(tac / core.price - 1) * 100:.1f}%) 추가 소요',
        size=10.5, color='ink2')

    y2 = 1.22 + 0.32 + 0.36 * 5 + 0.52
    rect(s, M, y2, left, 1.16, fill=K['panel'])
    txt(s, M + 0.20, y2 + 0.12, left - 0.40, 0.24, '연 수익률', size=11.5,
        bold=True, color='navy')
    for i, c_ in enumerate((core.gross_price, core.gross_price_deposit)):
        txt(s, M + 0.20, y2 + 0.42 + i * 0.30, left - 1.30, 0.26,
            c_.basis, size=11, color='ink2')
        txt(s, M + left - 1.30, y2 + 0.40 + i * 0.30, 1.10, 0.28,
            pct(c_.value), size=13, bold=True, align=PP_ALIGN.RIGHT)
    txt(s, M, y2 + 1.24, left, 0.24,
        '· 운영비 미제출 — "순수익률" 라벨을 쓰지 않음', size=10,
        color='ink2')

    x = M + left + 0.28
    w = RIGHT - x
    rows3, acc = [], {}
    for i, r in enumerate(core.ltv_rows):
        neg = r['monthly_net'] < 0
        rows3.append([f'{int(r["ltv"] * 100)}%',
                      eok(r['loan']) if r['loan'] else '—',
                      eok(r['equity']),
                      man(r['interest']) if r['interest'] else '—',
                      ('▼ ' if neg else '+') + man(r['monthly_net']),
                      pct(r['roe'])])
        if neg:
            acc[i] = 'red'
    table(s, x, 1.22, w,
          ['LTV', '대출금', '실투자금', '월 이자', '월 순현금', '자기자본 수익률'],
          rows3, [0.8, 1.1, 1.2, 1.0, 1.2, 1.4],
          size=10.5, row_h=0.38, align=['c', 'r', 'r', 'r', 'r', 'r'],
          accent_rows=acc)
    y3 = 1.22 + 0.32 + 0.38 * 3 + 0.24
    txt(s, x, y3, w, 0.22, '◇ 금리 4.5% 가정 (2026년 통상) · 이자만 · 운영비 미반영',
        size=10, color='mute')

    if core.negative_leverage:
        rect(s, x, y3 + 0.34, w, 1.06, fill=K['redbg'])
        txt(s, x + 0.18, y3 + 0.50, w - 0.36, 0.74,
            f'역레버리지 — 수익률 {pct(core.gross_price.value)} < 금리 4.5%\n'
            f'대출을 늘릴수록 자기자본 수익률이 낮아집니다. '
            f'무차입 {pct(core.roe_ceiling.value)}가 상한입니다.',
            size=11.5, bold=True, color='red', space=3, line_spacing=1.2)
    # 자금 계획 요약 — 좌하단
    ltv50 = {int(r['ltv'] * 100): r for r in core.ltv_rows}[50]
    yl = 1.22 + 0.32 + 0.36 * 5 + 0.52 + 1.16 + 0.52
    if yl + 1.30 <= BODY_Y:
        rect(s, M, yl, left, 1.28, fill=K['navy'])
        txt(s, M + 0.20, yl + 0.14, left - 0.40, 0.22,
            '· 자금 계획 (LTV 50% 기준)', size=10,
            color=RGBColor(0x9F, 0xAF, 0xC5))
        parts = [('총취득원가', eok(tac)), ('− 보증금', eok(core.deposit)),
                 ('− 대출', eok(ltv50['loan'])), ('= 실투자금', eok(ltv50['equity']))]
        for i, (lb, v) in enumerate(parts):
            yy = yl + 0.42 + (i // 2) * 0.40
            xx = M + 0.20 + (i % 2) * (left - 0.40) / 2
            txt(s, xx, yy, 1.15, 0.24, lb, size=10,
                color=RGBColor(0xA9, 0xB6, 0xC8))
            txt(s, xx + 1.15, yy - 0.02, 1.30, 0.28, v, size=13, bold=True,
                color='white' if i < 3 else RGBColor(0xF2, 0xC9, 0x8A))

    rect(s, x, y3 + 1.54, w, 0.94, fill=K['panel'])
    txt(s, x + 0.18, y3 + 1.68, w - 0.36, 0.68,
        '은행은 매매가가 아니라 감정가로 한도를 정합니다.\n'
        '잔금 일정 확정 전에 대출 승인을 받으시길 권합니다.',
        size=11, color='ink2', space=2, line_spacing=1.2)
    return s


def s09_market(prs, core, A):
    s = page(prs, 'Market Comparables', '가격 근거 — 인근 실거래',
             badge='국토부 실거래' if core.has_public else '확인 필요',
             badge_kind='ok' if core.has_public else 'warn')
    blk = core.blocks()
    if not blk['comps_table'][0]:
        rect(s, M, 1.30, CW, 1.30, fill=K['redbg'])
        txt(s, M + 0.24, 1.56, CW - 0.48, 0.80,
            '인근 실거래 비교사례 미확보\n'
            '시세 대비 위치와 목표 매각가를 산출하지 않았습니다.\n'
            + blk['comps_table'][1],
            size=13, bold=True, color='red', space=3, line_spacing=1.25)
        txt(s, M, 2.90, CW, 0.30,
            '· 단가 비교 시 분자와 분모의 기준을 맞춥니다\n'
            '· 연면적 단가를 토지 실거래와 비교하면 항상 몇 배로 나옵니다',
            size=11, color='ink2')
        return s
    comps = core.f('comps').value
    rows = [[n, d, f'{amt / 1e8:,.0f}억', f'{gfa:,.1f}㎡',
             f'{lpp:,}만', f'{gpp:,}만', fl] for n, d, amt, gfa, lpp, gpp, fl in comps]
    table(s, M, 1.22, CW,
          ['비교사례', '거래일', '거래금액', '연면적', '토지 평당가',
           '연면적 평당가', '규모'],
          rows, [3.4, 1.2, 1.3, 1.4, 1.5, 1.6, 1.2],
          size=10, row_h=0.32, align=['l', 'c', 'r', 'r', 'r', 'r', 'c'])
    y = 1.22 + 0.32 + 0.32 * len(rows) + 0.30
    llo, lhi = PD.comps_range(core.fixture_id, 'land')
    glo, ghi = PD.comps_range(core.fixture_id, 'gfa')
    lv = core.land_pyeong_price.value / MAN
    gv = core.gfa_pyeong_price.value / MAN
    rows2 = [['토지 평당가', f'{lv:,.0f}만원', f'{llo:,}만 ~ {lhi:,}만',
              f'{(lv / llo - 1) * 100:+.1f}%'],
             ['연면적 평당가', f'{gv:,.0f}만원', f'{glo:,}만 ~ {ghi:,}만',
              f'{(gv / glo - 1) * 100:+.1f}%']]
    table(s, M, y, 7.30, ['비교 기준', '본 자산', '인근 실거래 범위', '하단 대비'],
          rows2, [1.7, 1.5, 2.4, 1.4], size=10.5, row_h=0.36,
          align=['l', 'r', 'r', 'r'], accent_rows={0: 'sum'})

    x = M + 7.55
    w = RIGHT - x
    mult = core.land_price_multiple
    rect(s, x, y, w, 1.04, fill=K['panel'])
    txt(s, x + 0.18, y + 0.14, w - 0.36, 0.22, '· 매매가 ÷ 공시지가 총액',
        size=10.5, color='mute')
    txt(s, x + 0.18, y + 0.40, w - 0.36, 0.44, f'{mult.value:.2f}배',
        size=24, bold=True, color='navy')
    txt(s, x + 1.90, y + 0.52, w - 2.1, 0.40,
        f'· 공시지가 총액 {eok(core.land_price_total.value)}\n'
        f'· 매매가의 {core.land_price_total.value / core.price * 100:.1f}%',
        size=10, color='ink2', space=1)

    y2 = y + 0.32 + 0.36 * 2 + 0.30
    txt(s, M, y2, CW, 0.60,
        f'· 두 기준 모두 인근 하단보다 낮음\n'
        f'· 용적률 {core.f("farPct").value:.2f}%의 저밀 개발 — 연면적이 적어 토지 평당가가 낮게 형성\n'
        f'· 같은 사실이 잔여 용적률 {core.far_headroom.value:.1f}%p라는 개선 여력'
        if lv < llo else
        '· 토지·연면적 두 기준을 나눠 비교했습니다. 기준이 다르면 배수가 왜곡됩니다.',
        size=11, color='ink2', space=3, line_spacing=1.25)
    ycond = y2 + 0.86
    if ycond + 1.10 <= BODY_Y:
        rect(s, M, ycond, CW, 1.06, fill=K['panel'])
        txt(s, M + 0.24, ycond + 0.14, CW - 0.48, 0.24, '조사 조건', size=11,
            bold=True, color='navy')
        txt(s, M + 0.24, ycond + 0.44, CW - 0.48, 0.50,
            '· 조회 범위 — 인근 상업·업무용 실거래 직전 12개월 (국토부)\n'
            '· 실거래 API는 지번을 마스킹하고 집합건물 대지면적을 주지 않습니다\n'
            '· 물건 식별·접면·규모·대지면적은 중개인 보강값입니다',
            size=10, color='ink2', space=2, line_spacing=1.2)
    return s


def s10_risk(prs, core, A, confirmed):
    s = page(prs, 'Risk & Due Diligence', '리스크 및 확인사항',
             badge=f'확인 필요 {len(core.deficiencies)}건', badge_kind='warn')
    cw = (CW - 0.44) / 3
    cols = [('확인된 리스크', confirmed, 'red'),
            ('미확인 사항', core.deficiencies, 'gold'),
            ('매수인 확인 권고',
             ['등기부등본 — 근저당·권리 제한', '대출 감정가 — 잔금 전 확인',
              '임대차계약서 원본 — 최초 계약일·특약', '현장 — 누수·균열·설비'],
             'navy')]
    for i, (title, items, col) in enumerate(cols):
        x = M + i * (cw + 0.22)
        rect(s, x, 1.22, cw, 0.36, fill=K[col if col != 'gold' else 'gold'])
        txt(s, x + 0.16, 1.29, cw - 0.32, 0.24, title, size=11.5, bold=True,
            color='white')
        rect(s, x, 1.58, cw, 4.30, fill=None, line=K['line'])
        # 🔴 잘라내지 않습니다. 8개를 넘으면 마지막 칸에 남은 수를 적습니다.
        #    배지에는 11건이라 쓰고 8건만 그리면 결손 3건이 사라집니다
        #    (불변조건 13 위반). 실제로 그렇게 나가고 있었습니다.
        CAP = 8
        show = items if len(items) <= CAP else items[:CAP - 1]
        for j, it in enumerate(show):
            y = 1.74 + j * 0.52
            rect(s, x + 0.16, y + 0.08, 0.045, 0.16, fill=K[col])
            txt(s, x + 0.30, y - 0.02, cw - 0.46, 0.50, it, size=11,
                color='ink2', line_spacing=1.15)
        if len(items) > CAP:
            rest = len(items) - len(show)
            y = 1.74 + len(show) * 0.52
            rect(s, x + 0.16, y + 0.08, 0.045, 0.16, fill=K[col])
            txt(s, x + 0.30, y - 0.02, cw - 0.46, 0.50,
                f'외 {rest}건 — 모바일 IM 에 전량 표기', size=11,
                bold=True, color='ink2', line_spacing=1.15)
    txt(s, M, 6.06, CW, 0.50,
        '· 등기부등본 미제출 — 권리 제한 여부를 "없음"으로 판정하지 않았습니다\n'
        '· 확인사항 칸은 공개 단계에서도 가리지 않습니다',
        size=10.5, color='ink2', space=3)
    return s


def s11_photos(prs, core, A, keys, title, caps):
    s = page(prs, 'Photographs', title)
    if not keys:
        rect(s, M, 1.30, CW, 2.0, fill=K['panel'])
        txt(s, M, 2.10, CW, 0.40, '사진 미제출', size=16, bold=True,
            color='mute', align=PP_ALIGN.CENTER)
        txt(s, M, 2.56, CW, 0.30,
            '표준은 외부·내부 6장 이상을 필수로 둡니다. 촬영 후 보완이 필요합니다.',
            size=11, color='ink2', align=PP_ALIGN.CENTER)
        return s
    big = keys[0]
    photo_or_note(s, A, [big], M, 1.22, 7.10, 5.10)
    txt(s, M, 6.32, 7.10, 0.24, '· ' + caps[0], size=10, color='mute')
    x = M + 7.30
    w = RIGHT - x
    hh = (5.10 - 0.24 * 2) / 3
    for i, k in enumerate((keys[1:4] + [None, None, None])[:3]):
        y = 1.22 + i * (hh + 0.24)
        if k is None:
            continue
        photo_or_note(s, A, [k], x, y, w, hh)
        txt(s, x, y + hh + 0.03, w, 0.20,
            ('· ' + caps[i + 1]) if i + 1 < len(caps) else '',
            size=9.5, color='mute')
    return s


def sX_parcels(prs, core, A):
    """L10 — 필지 명세. 필지가 2개 이상일 때만 켭니다.

    🔴 단일 필지 물건에 이 면을 내면 "1필지" 라는 당연한 말을 한 면 씁니다.
       im.pages.yaml 의 min_parcels 가 그것을 막습니다.
    """
    L = core.land
    s = page(prs, 'Parcels & Exclusions', '필지 명세 · 제척',
             badge=f'{L.count}필지', badge_kind='ok')
    y0 = 1.22

    rows = []
    for pc in core.parcels:
        share = ('단독' if pc.ownership == 'sole'
                 else f'{pc.share_num}/{pc.share_den}')
        rows.append([pc.jibun, pc.jimok, f'{pc.area:,.1f}', share,
                     f'{pc.excluded_area:,.1f}' if pc.excluded_area else '—',
                     f'{pc.effective_area:,.1f}'])
    rows.append(['계', '', f'{L.ledger_area:,.1f}', '',
                 f'{L.excluded_area:,.1f}' if L.excluded_area else '—',
                 f'{L.effective_area:,.1f}'])
    table(s, M, y0, 7.10,
          ['지  번', '지목', '대장(㎡)', '지분', '제척(㎡)', '유효(㎡)'],
          rows, [1.9, 0.7, 1.2, 0.8, 1.1, 1.4], size=10, row_h=0.36,
          align=['l', 'c', 'r', 'c', 'r', 'r'],
          accent_rows={len(rows) - 1: 'sum'})

    y = y0 + 0.32 + 0.36 * len(rows) + 0.24
    for pc in core.parcels:
        for e in pc.exclusions:
            prov = PROV_KO.get(e.provenance, e.provenance)
            txt(s, M, y, 7.10, 0.26,
                f'· {pc.jibun} · {e.label} {e.area:,.1f}㎡ ●{prov}',
                size=11, color='ink2')
            y += 0.28
            txt(s, M + 0.18, y, 6.92, 0.26,
                '용적률 산정 제외' if e.affects_far else '용적률 산정 포함',
                size=11, color='mute')
            y += 0.28
            if e.note:
                # 45자를 넘지 않게 나눕니다 (정본 §4.7)
                for ln in _wrap45(e.note):
                    txt(s, M + 0.18, y, 6.92, 0.26, ln, size=11, color='mute')
                    y += 0.28

    # 우측 — 유효 면적이 무엇을 바꾸는가
    x = M + 7.30
    w = RIGHT - x
    rect(s, x, y0, w, 1.42, fill=K['panel'])
    txt(s, x + 0.18, y0 + 0.14, w - 0.36, 0.22, 'P01 유효 대지면적', size=11,
        color='mute')
    txt(s, x + 0.18, y0 + 0.40, w - 0.36, 0.52, f'{L.effective_area:,.1f}㎡',
        size=24, bold=True, color='ink')
    txt(s, x + 0.18, y0 + 0.98, w - 0.36, 0.30,
        f'대장 {L.ledger_area:,.1f}㎡ · 제척 {L.excluded_area:,.1f}㎡',
        size=11, color='ink2')

    y2 = y0 + 1.62
    rows2 = [['대장 기준', f'{L.ledger_far_pct:.1f}%'],
             ['P02 유효 기준', f'{L.effective_far_pct:.1f}%']]
    table(s, x, y2, w, ['용적률 기준', '값'], rows2, [1.7, 1.3],
          size=10, row_h=0.38, align=['l', 'r'], label_col=True,
          accent_rows={1: 'red'} if core.relief_cross else {1: 'sum'})

    rc = core.relief_cross
    y3 = y2 + 0.32 + 0.38 * 2 + 0.24
    if rc:
        rect(s, x, y3, w, 1.60, fill=K['redbg'])
        txt(s, x + 0.16, y3 + 0.14, w - 0.32, 0.24,
            '한시 완화 조건 이탈', size=11.5, bold=True, color='red')
        txt(s, x + 0.16, y3 + 0.44, w - 0.32, 1.02,
            f'{core.relief["name"]}\n'
            f'조건 {rc["threshold_pct"]:.0f}% 미만\n'
            f'대장 {rc["ledger_far_pct"]}% → 유효 {rc["effective_far_pct"]}%\n'
            f'{rc["action"]}',
            size=11, color='ink2', line_spacing=1.20)
    else:
        rect(s, x, y3, w, 0.90, fill=K['panel'])
        txt(s, x + 0.16, y3 + 0.28, w - 0.32, 0.40,
            '· 한시 완화 임계 이탈 없음', size=11, color='ink2')

    txt(s, M, 6.02, CW, 0.56,
        '· 제척 면적은 공공 API로 조회되지 않습니다\n'
        '· 토지이용계획도 판독값 — 관할 구청 확인 필요\n'
        '· 건폐율·용적률은 대장 대지 기준으로 고시됩니다',
        size=11, color='mute', line_spacing=1.20)
    return s


def sX_evidence(prs, core, A):
    """P7 전용 — 근거·검증 면. 무엇을 무엇으로 확인했는가."""
    s = page(prs, 'Evidence & Verification', '근거 · 검증',
             badge='공부 교차검증', badge_kind='ok')
    chip_band(s, core, 1.06)
    y0 = 1.46
    xr = PD.crosscheck(core.fixture_id)
    def verdict(x) -> str:
        """판정 문구.

        🔴 부등식 검사(계약면적 합 ≤ 연면적)에 "일치 (11.13%)" 를 쓰면
           읽는 사람이 11.13% 어긋난 것으로 봅니다. 그 11.13% 는 공용부
           비율이지 오차가 아닙니다. 검사 성격에 맞는 말을 씁니다.
        """
        if '≤' in x.label:
            return '충족' if x.ok else '초과'
        return ('일치' if x.ok else '불일치') + f' ({x.gap_pct:.2f}%)'

    rows = [[x.code, x.label,
             f'{x.expected:,.2f}', f'{x.actual:,.2f}', verdict(x)]
            for x in xr]
    table(s, M, y0, 7.30, ['코드', '검증 내용', '기대', '실제', '판정'],
          rows, [0.7, 3.0, 1.4, 1.4, 1.3], size=10, row_h=0.36,
          align=['c', 'l', 'r', 'r', 'c'])

    x = M + 7.55
    w = RIGHT - x
    Lg, Pg, short = core.resolution_pair()
    rect(s, x, y0, w, 1.30, fill=K['panel'])
    txt(s, x + 0.20, y0 + 0.14, w - 0.40, 0.22, '· 자료 해상도', size=10,
        color='mute')
    txt(s, x + 0.20, y0 + 0.40, w - 0.40, 0.44, f'{Lg} · {Pg}', size=24,
        bold=True, color='navy')
    txt(s, x + 0.20, y0 + 0.92, w - 0.40, 0.26,
        '· 임대차 L축 / 물건자료 P축', size=10.5, color='ink2')

    blk = core.blocks()
    off = [(k, v[1]) for k, v in blk.items() if not v[0]]
    y1 = y0 + 0.32 + 0.36 * len(rows) + 0.34
    txt(s, M, y1, CW, 0.28, f'잠긴 블록 {len(off)}종 — 사유', size=13,
        bold=True, color='navy')
    y2 = y1 + 0.38
    cw2 = CW / 2 - 0.16
    for i, (k, why) in enumerate(off[:10]):
        b = ISPEC.BLOCK_BY_KEY[k]
        xx = M + (i % 2) * (cw2 + 0.32)
        yy = y2 + (i // 2) * 0.32
        if yy + 0.30 > BODY_Y:
            break
        txt(s, xx, yy, 2.30, 0.26, '· ' + b.label, size=10.5, bold=True)
        txt(s, xx + 2.40, yy, cw2 - 2.40, 0.26, why[:34], size=11,
            color='ink2')
    return s


def sX_landvalue(prs, core, A):
    """P8 전용 — 토지가치 면. 한국 매수자 판단축 2·3."""
    s = page(prs, 'Land Value', '토지 가치 · 개발 여력',
             badge='공부 · 공시지가', badge_kind='ok')
    f = core.f
    cw2 = (CW - 0.44) / 3
    cards = [('토지 평당가',
              f'{core.land_pyeong_price.value / MAN:,.0f}만원',
              f'매매가 ÷ 대지 {core.land_sqm.value / PYEONG:,.1f}평'),
             ('공시지가 배수', f'{core.land_price_multiple.value:.2f}배',
              f'공시지가 총액 {eok(core.land_price_total.value)}'),
             ('잔여 용적률', f'{core.far_headroom.value:.1f}%p',
              f'현행 {f("farPct").value:.2f}% / 상한 {f("farLimit").value:.0f}%')]
    for i, (lb, v, sub) in enumerate(cards):
        x = M + i * (cw2 + 0.22)
        rect(s, x, 1.30, cw2, 1.50, fill=K['panel'])
        txt(s, x + 0.20, 1.46, cw2 - 0.40, 0.24, lb, size=11, color='ink2')
        txt(s, x + 0.20, 1.74, cw2 - 0.40, 0.50, v, size=26, bold=True,
            color='gold')
        txt(s, x + 0.20, 2.32, cw2 - 0.40, 0.34, '· ' + sub, size=10.5,
            color='mute')

    # 용적률 게이지
    gy = 3.06
    txt(s, M, gy, CW, 0.28, '용적률 여력', size=13, bold=True, color='navy')
    cur, lim = f('farPct').value, f('farLimit').value
    gw = CW
    rect(s, M, gy + 0.40, gw, 0.44, fill=K['panel'], line=K['line'])
    rect(s, M, gy + 0.40, gw * min(cur / lim, 1.0), 0.44, fill=K['navy'])
    txt(s, M + 0.16, gy + 0.51, 4.0, 0.24, f'현행 {cur:.2f}%', size=12,
        bold=True, color='white')
    txt(s, RIGHT - 4.0 - 0.16, gy + 0.51, 4.0, 0.24, f'법정 상한 {lim:.0f}%',
        size=12, color='ink2', align=PP_ALIGN.RIGHT)

    ry = gy + 1.02
    rows = [['용도지역', f('zoning').value, f('zoning').source],
            ['법정 건폐율 / 용적률',
             f'{f("bcrLimit").value:.0f}% / {lim:.0f}%', f('farLimit').source],
            ['현행 건폐율 / 용적률',
             f'{f("bcrPct").value:.2f}% / {cur:.2f}%', '건축물대장'],
            ['㎡당 공시지가', f'{f("landPriceSqm").value:,}원',
             f('landPriceSqm').source],
            ['토지 공시지가 총액', eok(core.land_price_total.value),
             '대지면적 × ㎡당 공시지가']]
    table(s, M, ry, CW, ['구  분', '값', '출처'], rows,
          [2.6, 4.2, 5.3], size=10.5, row_h=0.34, align=['c', 'l', 'l'],
          label_col=True)
    yb = ry + 0.32 + 0.34 * len(rows) + 0.20
    if yb + 0.30 <= BODY_Y:
        txt(s, M, yb, CW, 0.28,
            '· 증축 규모는 가로구역 최고높이·주차 기준 검토 후 산출'
            if core.far_headroom.value > 20
            else '· 법정 상한까지 개발 완료 — 증축 여력 사실상 없음',
            size=11, color='ink2')
    return s


def s12_terms(prs, core, A, steps):
    s = page(prs, 'Terms & Next Steps', '거래 조건 · 다음 단계')
    left = 5.60
    rows = [['매매 희망가', eok(core.price, 0) + '원'],
            ['승계 보증금', eok(core.deposit) + '원'],
            ['월 임대료', man(core.monthly_rent) + '원 (VAT 별도)'],
            ['잔금 일정', '협의'],
            ['대출 승계', '확인 필요'],
            ['소유 형태', '층별 구분등기 · 소유자 2인'
             if core.fixture_id == 'dangsan' else '확인 필요']]
    table(s, M, 1.22, left, ['구  분', '내  용'], rows, [1.7, 3.4],
          size=10.5, row_h=0.40, align=['c', 'l'], label_col=True,
          accent_rows={0: 'red'})

    x = M + left + 0.28
    w = RIGHT - x
    txt(s, x, 1.22, w, 0.28, '진행 절차', size=13, bold=True, color='navy')
    for i, st in enumerate(steps):
        y = 1.62 + i * 0.50
        rect(s, x, y, 0.30, 0.30, fill=K['navy'])
        txt(s, x, y + 0.055, 0.30, 0.22, str(i + 1), size=11, bold=True,
            color='white', align=PP_ALIGN.CENTER)
        txt(s, x + 0.42, y + 0.03, w - 0.42, 0.42, st, size=11.5, color='ink2')

    y = 1.22 + 0.32 + 0.40 * 6 + 0.34
    rect(s, M, y, left, 1.10, fill=K['navy'])
    txt(s, M + 0.22, y + 0.16, left - 0.44, 0.24, '문의', size=10,
        color=RGBColor(0x9F, 0xAF, 0xC5))
    txt(s, M + 0.22, y + 0.42, left - 0.44, 0.30, '제이에스부동산중개(주)',
        size=14, bold=True, color='white')
    txt(s, M + 0.22, y + 0.76, left - 0.44, 0.24,
        '담당 중개인  ·  02-000-0000', size=11,
        color=RGBColor(0xC9, 0xD4, 0xE2))
    # 자료가 채워지면 무엇이 열리는가 — 잠긴 블록에서 역산합니다
    blk = core.blocks()
    want: dict[str, list[str]] = {}
    for b in ISPEC.BLOCKS:
        if blk[b.key][0]:
            continue
        for k in b.needs:
            if k in ISPEC.FIELD_BY_KEY:
                want.setdefault(k, []).append(b.label)
    top = sorted(want.items(), key=lambda kv: -len(kv[1]))[:3]
    if top:
        yb = y + 1.28
        rect(s, x, yb - 0.06, w, 1.44, fill=K['panel'])
        txt(s, x + 0.20, yb + 0.08, w - 0.40, 0.24,
            '· 자료를 넣으면 열리는 내용', size=10.5, color='mute')
        for i, (k, labels) in enumerate(top):
            txt(s, x + 0.20, yb + 0.36 + i * 0.34, 2.30, 0.28,
                '· ' + ISPEC.label(k), size=11, bold=True, color='navy')
            txt(s, x + 2.55, yb + 0.36 + i * 0.34, w - 2.75, 0.28,
                '→ ' + ' · '.join(labels[:2]), size=10.5, color='ink2')

    txt(s, M, y + 1.24, CW, 0.36,
        '· 문의 시 물건명·임차인 상호는 가려서 보내주십시오\n'
        '· 본 자료의 수치는 제출된 임대 현황과 공부에서 산출했습니다',
        size=9.5, color='mute')
    return s


# ── 물건별 카피 (개조식) ───────────────────────────────────────────────
COPY = {
    'dangsan': {
        'alias': '근생빌딩',
        'upside_line': '· 자가사용 2구획 임대 전환\n'
                       '· 11년간 인상 이력 없는 계약의 재계약 조정\n'
                       '· 금액은 시장 임차료 자료 확보 후 산출',
        'tags_R2': ['당산역 380m', '준공업지역', '의료 앵커'],
        'tags_R1': ['임대 6구획', '공실 0'],
        'points_R2': [
            '당산역(2·9호선) 380m · 도보 5분 — 여의도 1정거장',
            '준공업지역 용적률 221.8% / 상한 400% — 잔여 178.2%p',
            '토지 평당 7,492만원 — 인근 실거래 하단 9,300만원 대비 낮음',
            '의원·약국 등 의료 앵커 6구획 임차 · 공실 0',
            '공시지가 총액 60.13억 — 매매가의 52.3%가 땅값',
        ],
        'points_R1': [
            '임대 6구획 · 자가사용 2구획 · 공실 0',
            '월 임대료 1,946만원 · 보증금 2.90억원 (원장 8행 합계)',
            '연면적 미확정 — 계 행 1,141.15㎡ vs 층별 합 1,441.15㎡',
            '용도지역·공시지가·인근 실거래 미조회',
            '확인 필요 7건 — 공부 결합 시 대부분 해소',
        ],
        'bullets_R2': [
            '지하 1층~지상 5층 · 코너 대로변 · 배후 아파트 5,000세대',
            '2002년 준공 · 위반건축물 없음 · 철근콘크리트',
            '층별 구분등기 · 소유자 2인 — 매각 동의 선결',
        ],
        'bullets_R1': [
            '지하 1층~지상 5층 · 8개 구획',
            '공부 미조회 — 규모·구조·준공 확인 필요',
        ],
        'loc_notes': [
            '국회대로·올림픽대로 400m — 여의도·강남 접근 용이',
            '배후 아파트 5,000세대 · 영등포구청·세무서 인근',
            '상권 — 의료·보건 24.5% · 음식 38.2% (소상공인 상권정보)',
        ],
        'upside': [
            '자가사용 2구획 임대 전환 여지\n— 시장 임차료 자료 없어 금액 미산출',
            '11년간 임대료 인상 이력 없음 (원장 비고) — 재계약 시 조정 여지',
            '◇ 잔여 용적률 178.2%p — 증축 규모는 높이·주차 기준 검토 후 산출',
        ],
        'risk': [
            '4F 주류판매 계약 만료 경과 (2025-04-30)',
            '역레버리지 — LTV 50%에서 월 −210만원',
            '1F·2F 통합계약 — 분리 임대 시 재계약',
            '소유자 2인 구분등기 — 전원 동의 필요',
        ],
        'ext': ['hero', 'front', 'ext1', 'ext2'],
        'ext_caps': ['코너 전경 — 횡단보도 접면', '정면', '측면 진입부', '가로 전경'],
        'inr': ['in_cafe', 'in_lobby', 'in_shop', 'in_ev'],
        'inr_caps': ['지하 1층 카페 (자가사용)', '1층 로비', '1층 약국', '승강기 홀'],
    },
    'yangpyeong': {
        'alias': '오피스빌딩',
        'upside_line': '· 지하 1층 422.25㎡(127.7평) 공실 임대\n'
                       '· 금액은 비교 임대사례 확보 후 산출',
        'tags_R2': ['선유도역 95m', '2018년 준공', '용적률 398.8%'],
        'tags_R1': ['임대 11구획', '공실 1'],
        'points_R2': [
            '선유도역(9호선) 95m · 도보 1분 — 4번 출구 대로변 직결',
            '2018년 준공 · 지상 10층 · 용적률 398.8% (상한 400%)',
            '토지 평당 {land_pyeong} — 인근 1.38억~1.85억 범위 내',
            'IT·법률·디자인 등 11개 법인 임차 · 지하 1층 공실',
            '기계식 22 + 자주식 1 · 15인승 승강기 · 개별 EHP',
        ],
        'points_R1': [
            '임대 11구획 · 공실 1구획 (지하 1층)',
            '월 임대료 4,657만원 · 관리비 576만원 (원장 12행 합계)',
            '표지 요약 5,017만원과 원장 합계 4,657만원 불일치',
            '첨부 토지이용계획확인원이 타 자치구 필지',
            '임대면적·용도지역·공시지가 미확보',
        ],
        'bullets_R2': [
            '3필지 합지 518.70㎡ · 건폐율 58.40% — 최적 설계',
            '법정 상한 400% 대비 398.80% — 개발 완료 자산',
            '위반건축물 없음 · 철근콘크리트 · 지하 1층~지상 10층',
        ],
        'bullets_R1': [
            '지하 1층~지상 10층 · 12개 구획',
            '공부 미조회 — 대지·연면적·용적률 확인 필요',
        ],
        'loc_notes': [
            '양평로·올림픽대로·서부간선 300m — 여의도·마포·DMC 연결',
            '선유도 IT·지식산업 클러스터 · 선유도공원 인접',
            '상권 — IT·전문서비스 54.2% (소상공인 상권정보)',
        ],
        'upside': [
            '지하 1층 422.25㎡(127.7평) 공실\n— 임대 전환 시 수입 증가\n— 비교 임대사례 없어 금액 미산출',
            '11개 법인 임차 · 업종 분산 — 만기 자료만으로는 공실 위험 계량 불가',
            '◇ 용적률 잔여 1.2%p — 증축 여력 사실상 없음',
        ],
        'risk': [
            '표지 요약과 원장 합계 월 360만원 불일치 (연 4,320만원)',
            '역레버리지 — LTV 50%에서 월 −30만원',
            '지하 1층 공실 — 면적 기준 공실률 16.95%',
            '첨부 공부 오첨부 (R1판) — 정본 필지 교체 필요',
        ],
        'ext': ['hero', 'context'],
        'ext_caps': ['건물 전경 — 선유도역 4번 출구 대로변', '배후 주거·업무 밀집'],
        'inr': [],
        'inr_caps': [],
    },
}

COPY['multiparcel'] = {
    'alias': '근생빌딩',
    'upside_line': '· 지하 1층 327.8㎡(99.2평) 공실 임대\n'
                   '· 금액은 비교 임대사례 확보 후 산출',
    'tags_R2': ['잠원역 450m', '2필지', '제척 12.5㎡'],
    'tags_R1': ['임대 5구획', '공실 1'],
    'points_R2': [
        '잠원역(3호선) 450m · 도보 6분 — 강남대로 접근 600m',
        '2필지 616.1㎡ — 제척 12.5㎡ 반영 시 유효 603.6㎡',
        '토지 평당 {land_pyeong} — 인근 1.31억~1.62억 범위 내',
        '약국·의원·학원 등 5개 임차 · 지하 1층 공실',
        '한시 완화 250% 조건 — 유효 기준 254.1%로 이탈',
    ],
    'points_R1': [
        '임대 5구획 · 공실 1구획 (지하 1층)',
        '보증금 3.20억 · 월세 2,140만원',
        '공부 미조회 — 대지·용적률 확인 필요',
    ],
    'bullets_R2': ['2필지 · 대장 616.1㎡ · 유효 603.6㎡',
                   '제척 12.5㎡ — 9M 도로 확폭 계획 저촉',
                   '지상 5층 · 지하 1층 · 2009년 준공'],
    'bullets_R1': ['임대 5구획 · 공실 1구획', '공부 조회 후 확정'],
    'loc_notes': ['잠원역 3호선 450m · 강남대로 600m',
                  '배후 아파트 밀집 · 강남 업무지구 연계',
                  '상권 — 의료 21.4% · 음식 33.8% (소상공인 상권정보)'],
    'upside': ['지하 1층 327.8㎡(99.2평) 공실',
               '금액은 비교 임대사례 확보 후 산출',
               '제척 확정 후 증축 여력 재검토'],
    'risk': ['제척 12.5㎡ — 한시 완화 250% 조건 이탈',
             '지하 1층 공실 — 면적 기준 공실률 18.32%',
             '제척은 도면 판독값 — 관할 구청 확인 필요'],
    'ext': [], 'ext_caps': [], 'inr': [], 'inr_caps': [],
}

STEPS = ['관심 표명 · 담당자 연락', 'NDA 체결 · 임대차계약서 열람',
         '현장 실사 일정 조율', 'LOI 제출 · 가격 협의',
         '법률·세무 실사', '매매계약 체결 · 잔금']


def build(fid: str, edition: str, preset: str = 'jsre_field_navy') -> Presentation:
    apply_preset(PS.PRESETS[preset])
    core = im_core.load(fid, edition=edition)
    cp = COPY[fid]
    A = {p.stem: p for p in (HERE / 'assets' / fid).glob('*.jpg')}
    R = edition
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)

    s01_cover(prs, core, A)
    s02_points(prs, core, A, cp[f'points_{R}'])
    s03_overview(prs, core, A, cp[f'tags_{R}'], cp[f'bullets_{R}'])
    s04_location(prs, core, A, cp['loc_notes'] if R == 'R2' else
                 ['접면 도로·배후 시설 현장 확인 필요',
                  '역 거리 미확보 — 도보 시간 임의 기재하지 않음'])
    s05_land(prs, core, A)
    s06_rentroll(prs, core, A)
    s07_lease2(prs, core, A, cp['upside'])
    s08_invest(prs, core, A)
    s09_market(prs, core, A)
    s10_risk(prs, core, A, cp['risk'])
    if PRESET.switches.get('crosscheckPage') and core.has_public:
        sX_evidence(prs, core, A)
    if PRESET.switches.get('farGaugePage') and core.has_public:
        sX_landvalue(prs, core, A)
    if cp['ext']:
        s11_photos(prs, core, A, cp['ext'], '건물 사진 — 외부', cp['ext_caps'])
    if cp['inr']:
        s11_photos(prs, core, A, cp['inr'], '건물 사진 — 내부', cp['inr_caps'])
    else:
        s11_photos(prs, core, A, [], '건물 사진 — 내부', [])
    s12_terms(prs, core, A, STEPS)
    footer(prs, core)
    cap = PRESET.switches.get('maxSlides', 16)
    if len(prs.slides) > cap:
        raise Overflow(f'{len(prs.slides)}장 — 프리셋 {PRESET.key} 상한 {cap}장 초과')
    return prs


def main() -> int:
    for fid, name in (('dangsan', '당산동'), ('yangpyeong', '양평동')):
        for ed in ('R1', 'R2'):
            prs = build(fid, ed, 'jsre_field_navy')
            p = HERE / f'{name}_PPTX_IM_KR_{ed}.pptx'
            prs.save(str(p))
            print(f'{p.name:<32} {len(prs.slides):>2}장 · '
                  f'{p.stat().st_size / 1024:>6,.0f} KB')
    # 프리셋 데모 — 당산동 R2로 신규 3종 비교
    for key in PS.PRESETS:
        prs = build('dangsan', 'R2', key)
        p = HERE / f'데모_당산동_{key}.pptx'
        prs.save(str(p))
        print(f'{p.name:<40} {len(prs.slides):>2}장 · '
              f'{p.stat().st_size / 1024:>6,.0f} KB')
    return 0


if __name__ == '__main__':
    try:
        sys.exit(main())
    except Overflow as e:
        print(f'경계 위반: {e}', file=sys.stderr)
        sys.exit(1)
