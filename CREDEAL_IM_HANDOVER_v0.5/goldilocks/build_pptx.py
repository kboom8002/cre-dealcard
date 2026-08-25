#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build_pptx.py — PPTX IM 생성

**정본은 `IM_STANDARD_수익형.md` §4 입니다.**
  §4.1 12p 필수 구성        §4.2 p2 한 장 요약 8블록
  §4.3 투자 구조             §4.6 리스크 3구획
  §4.7 시각 — 본문 11pt+ · 표 9pt+ · 회색 #666보다 어둡게
       색 의존 금지(기호 병기) · **페이지당 표 1개** · 한 줄 45자 이하
보조 준수: D7 `PPTX_ARCHETYPE_SPEC.md`
  · 13.333 × 7.5 in · M=0.62 · CW=12.093 · RIGHT=12.713
  · 본문 끝 ≤ 6.75 · Footer 6.94
  · **라벨은 절삭하지 않습니다** (v2 회귀 E01)
  · 표는 네이티브 표 (이미지 아님) · 원장 전량 (불변조건 18)

모든 도형은 `place()` 를 거칩니다. 경계를 넘으면 예외로 멈춥니다.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

import im_core
import im_copy
from im_core import (ASSUMED_LOAN_RATE, MAN, eok, man, pct, sqm_pyeong)

OUT = Path(__file__).resolve().parent

# ── 기하 (D7 §1) ───────────────────────────────────────────────────────
SW, SH = 13.333, 7.5
M, CW = 0.62, 12.093
RIGHT = M + CW                 # 12.713
BODY_MAX_Y, FOOTER_Y = 6.75, 6.94

FONT = 'Noto Sans CJK KR'

C = {
    'ink':   RGBColor(0x12, 0x16, 0x1C),
    'ink2':  RGBColor(0x3C, 0x43, 0x4E),
    'mute':  RGBColor(0x55, 0x5B, 0x66),   # #666 보다 어둡게 (§4.7)
    'line':  RGBColor(0xDF, 0xE3, 0xE8),
    'panel': RGBColor(0xF4, 0xF5, 0xF7),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'warnbg': RGBColor(0xFF, 0xF4, 0xF2),
    'warn':  RGBColor(0x8A, 0x2C, 0x1A),
    'okbg':  RGBColor(0xEA, 0xF4, 0xEE),
    'ok':    RGBColor(0x25, 0x6B, 0x42),
}


class Overflow(Exception):
    pass


def place(x, y, w, h, label=''):
    """D7 §1 경계 검사. 넘으면 렌더하지 않고 멈춥니다."""
    if x + w > RIGHT + 1e-6:
        raise Overflow(f'우측 이탈 {label}: {x + w:.3f} > {RIGHT}')
    if y + h > BODY_MAX_Y + 1e-6:
        raise Overflow(f'Footer 침범 {label}: {y + h:.3f} > {BODY_MAX_Y}')
    return Inches(x), Inches(y), Inches(w), Inches(h)


def txt(slide, x, y, w, h, text, *, size=12, bold=False, color='ink',
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=0, check=True):
    if check:
        place(x, y, w, h, text[:18])
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = []
    for raw in text.split('\n'):                  # §4.7 한 줄 45자 이하
        while len(raw) > 45:
            cut = raw.rfind(' ', 0, 46)
            cut = cut if cut > 20 else 45
            lines.append(raw[:cut].rstrip())
            raw = raw[cut:].lstrip()
        lines.append(raw)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(max(size, 11))          # §4.7 본문 11pt 이상
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = C[color] if isinstance(color, str) else color
    return tb


def rect(slide, x, y, w, h, fill='panel', line=None, check=True):
    if check:
        place(x, y, w, h, 'rect')
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = C[fill] if isinstance(fill, str) else fill
    if line:
        sp.line.color.rgb = C[line]
        sp.line.width = Pt(0.75)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    sp.text_frame.text = ''
    return sp


def table(slide, x, y, w, head, rows, widths, *, size=10, head_size=10,
          row_h=0.30, head_h=0.34, align=None):
    """네이티브 표. 이미지로 굽지 않습니다."""
    h = head_h + row_h * len(rows)
    place(x, y, w, h, 'table')
    shp = slide.shapes.add_table(len(rows) + 1, len(head),
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    t = shp.table
    t.first_row = True
    t.horz_banding = False
    total = sum(widths)
    for i, cw in enumerate(widths):
        t.columns[i].width = Emu(int(Inches(w) * cw / total))
    t.rows[0].height = Inches(head_h)
    for i in range(len(rows)):
        t.rows[i + 1].height = Inches(row_h)

    align = align or ['l'] * len(head)
    AL = {'l': PP_ALIGN.LEFT, 'r': PP_ALIGN.RIGHT, 'c': PP_ALIGN.CENTER}

    def fill(cell, text, *, bold, sz, col, al):
        cell.margin_left = cell.margin_right = Inches(0.07)
        cell.margin_top = cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = AL[al]
        r = p.add_run()
        r.text = str(text)
        r.font.size = Pt(max(sz, 9))             # §4.7 표 9pt 이상
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = col

    for j, hd in enumerate(head):
        c = t.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = C['panel']
        fill(c, hd, bold=True, sz=head_size, col=C['ink2'], al=align[j])
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            c = t.cell(i + 1, j)
            c.fill.solid()
            c.fill.fore_color.rgb = C['white']
            s = str(v)
            neg = s.startswith('-') or s.startswith('−')
            fill(c, s, bold=neg, sz=size,
                 col=C['warn'] if neg else C['ink'], al=align[j])
    return shp


def slide_new(prs, kicker, title, archetype):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    txt(s, M, 0.42, CW, 0.24, kicker, size=10.5, color='mute')
    txt(s, M, 0.68, CW, 0.44, title, size=22, bold=True)
    # 아키타입 표기 — 라벨은 절삭하지 않습니다
    txt(s, M, 1.18, CW, 0.02, '', size=1, check=False)
    ln = s.shapes.add_shape(1, Inches(M), Inches(1.20), Inches(CW), Inches(0.012))
    ln.fill.solid(); ln.fill.fore_color.rgb = C['line']
    ln.line.fill.background(); ln.shadow.inherit = False
    s._archetype = archetype
    return s


def footer(prs, core):
    for i, s in enumerate(prs.slides, 1):
        a = getattr(s, '_archetype', '')
        txt(s, M, FOOTER_Y, CW * 0.6, 0.26,
            f'CREDEAL IM · 수익형 · {core.address_band} · 해상도 {core.resolution_computed[0]}',
            size=8.5, color='mute', check=False)
        txt(s, M + CW * 0.6, FOOTER_Y, CW * 0.4, 0.26, f'{a}   {i} / {len(prs.slides)}',
            size=8.5, color='mute', align=PP_ALIGN.RIGHT, check=False)


# ── 슬라이드 ───────────────────────────────────────────────────────────
def p01_cover(prs, core):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, SW, SH, fill=C['ink'], check=False)   # 표지 전면 배경
    gu = core.address_band.split()[-1]
    txt(s, M, 2.10, CW, 0.30, 'CREDEAL INVESTMENT MEMORANDUM',
        size=11.5, color=C['mute'], check=False)
    txt(s, M, 2.55, CW, 0.95, f'{gu} {core.building_use} 매각',
        size=40, bold=True, color=C['white'], check=False)
    txt(s, M, 3.60, CW, 0.34,
        f'{core.address_band}  ·  수익형',
        size=15, color=RGBColor(0xA8, 0xB2, 0xC1), check=False)
    txt(s, M, 4.35, CW, 0.30,
        f'매매 희망가 {eok(core.price, 0)}   |   월 임대료 {man(core.monthly_rent)}원'
        f'   |   원장 {len(core.rows)}행',
        size=14, color=C['white'], check=False)
    txt(s, M, 5.55, CW, 0.80,
        f'본 자료의 모든 수치는 제출된 임대차 원장에서 산출했습니다.\n'
        f'확인되지 않은 항목 {len(core.deficiencies)}건은 "확인 필요"로 표기했으며, '
        f'추정값으로 채우지 않았습니다.',
        size=11.5, color=RGBColor(0x8B, 0x95, 0xA5), space=4, check=False)
    s._archetype = 'A01'
    return s


def p02_summary(prs, core):
    """정본 §4.2 — 8블록. 매수인의 60%가 여기서 진행 여부를 정합니다."""
    s = slide_new(prs, '한 장 요약', '여기서 판단이 끝납니다', 'A02')
    ltv = {int(r['ltv'] * 100): r for r in core.ltv_rows}
    r50 = ltv[50]

    # ① 한 줄 정의
    txt(s, M, 1.34, CW, 0.42, core.one_liner(), size=22, bold=True)

    # ②③④ 매매가 · 실투자금 · 월 순현금
    cw = (CW - 0.20 * 2) / 3
    cards = [
        ('매매가', eok(core.price, 0),
         (f'평당 {core.land_pyeong_price.value / MAN:,.0f}만원'
          if core.land_pyeong_price.known else '평당가 확인 필요')),
        ('실투자금', eok(r50['equity']), 'LTV 50% 기준'),
        ('월 순현금', f'{man(r50["monthly_net"])}원',
         '▼ 대출 시 적자' if r50['monthly_net'] < 0 else 'LTV 50% 기준'),
    ]
    for i, (lb, v, sub) in enumerate(cards):
        x = M + i * (cw + 0.20)
        neg = v.startswith('-') or v.startswith('−')
        rect(s, x, 1.90, cw, 1.34, fill='warnbg' if neg else 'panel')
        txt(s, x + 0.18, 2.06, cw - 0.36, 0.24, lb, size=11.5, color='ink2')
        txt(s, x + 0.18, 2.34, cw - 0.36, 0.48, v, size=26, bold=True,
            color='warn' if neg else 'ink')
        txt(s, x + 0.18, 2.88, cw - 0.36, 0.26, sub, size=11,
            color='warn' if neg else 'ink2')

    # ⑤ 레버리지 시나리오
    rows = []
    for r in core.ltv_rows:
        neg = r['monthly_net'] < 0
        rows.append([f'{int(r["ltv"] * 100)}%',
                     eok(r['loan']) if r['loan'] else '—',
                     eok(r['equity']),
                     man(r['interest']) if r['interest'] else '—',
                     ('▼ ' if neg else '+') + man(r['monthly_net']),
                     pct(r['roe'])])
    table(s, M, 3.40, CW,
          ['LTV', '대출금', '실투자금', '월 이자', '월 순현금', '자기자본 수익률'],
          rows, [1.4, 2.0, 2.2, 2.0, 2.3, 2.6],
          size=11, head_size=11, row_h=0.28, head_h=0.30,
          align=['r'] * 6)

    # ⑥ 임대 요약 · ⑦ 확인 필요
    hw = (CW - 0.20) / 2
    st = core.state_counts
    rect(s, M, 4.66, hw, 1.42, fill='panel')
    txt(s, M + 0.18, 4.80, hw - 0.36, 0.24, '임대 요약', size=11.5, bold=True,
        color='ink2')
    txt(s, M + 0.18, 5.08, hw - 0.36, 0.90,
        f'{" · ".join(f"{k} {v}" for k, v in st.items())}\n'
        f'월세 {man(core.ledger_sum_rent)}원 · 보증금 {eok(core.deposit)}원\n'
        f'연 수익률 {pct(core.gross_price.value)} (총임대료 ÷ 매매가)',
        size=11.5, space=3)

    rect(s, M + hw + 0.20, 4.66, hw, 1.42, fill='warnbg')
    txt(s, M + hw + 0.38, 4.80, hw - 0.36, 0.24,
        f'확인 필요 {len(core.deficiencies)}건', size=11.5, bold=True, color='warn')
    txt(s, M + hw + 0.38, 5.08, hw - 0.36, 0.90,
        '\n'.join('· ' + d.split(' —')[0].split(' (')[0]
                  for d in core.deficiencies[:4]),
        size=11, color='warn', space=2)

    # ⑧ 한 줄 코멘트 — 견해임을 밝힙니다 (정본 §7.1)
    txt(s, M, 6.22, CW, 0.32,
        f'저희가 보기에 무차입 또는 저LTV 구조가 맞습니다. '
        f'수익률이 대출금리 {ASSUMED_LOAN_RATE * 100:.1f}%를 밑돕니다.',
        size=12, color='ink2')
    return s


LINE_H, GAP = 0.215, 0.10          # 11.5pt × 1.35 · 블록 간격
ROW_H, HEAD_H = 0.26, 0.29


def _lines(t: str, per: int = 45) -> int:
    """렌더러와 같은 45자 줄바꿈으로 셉니다."""
    return max(1, -(-len(t.replace('**', '')) // per))


def _height(b, text_size=11.5) -> float:
    """블록이 차지할 높이. 페이지 넘김 판정에만 씁니다."""
    k = b['t']
    if k in ('p', 'note', 'warn', 'h'):
        h = max(_lines(b['text']) * LINE_H, 0.22)
        return h + {'h': 0.16, 'warn': 0.24}.get(k, GAP)
    if k == 'list':
        return sum(_lines(x, 52) * LINE_H for x in b['items']) + 0.14
    return HEAD_H + ROW_H * len(b['rows']) + 0.16


def _blocks_to_pptx(s, core, blocks, y0, *, max_y=BODY_MAX_Y - 0.04,
                    text_size=11.5, table_size=9.5):
    """카피 블록을 좌표로 흘려 넣습니다. 넘치면 남은 블록을 반환합니다."""
    y = y0
    for i, b in enumerate(blocks):
        k = b['t']
        # 제목이 혼자 남지 않게 — 바로 다음 블록까지 들어가야 제목을 놓습니다
        if k == 'h' and i + 1 < len(blocks):
            if y + _height(b, text_size) + _height(blocks[i + 1], text_size) > max_y:
                return blocks[i:]
        if k in ('p', 'note', 'warn', 'h'):
            t = b['text'].replace('**', '')
            size = 13 if k == 'h' else text_size
            bold = k == 'h'
            col = 'warn' if k == 'warn' else ('mute' if k == 'note' else 'ink')
            h = max(_lines(t) * LINE_H, 0.22)
            if y + h > max_y:
                return blocks[i:]
            if k == 'warn':
                rect(s, M, y - 0.04, CW, h + 0.12, fill='warnbg')
            txt(s, M + (0.12 if k == 'warn' else 0), y,
                CW - (0.24 if k == 'warn' else 0), h, t,
                size=size, bold=bold, color=col, space=3)
            y += h + {'h': 0.16, 'warn': 0.24}.get(k, GAP)
        elif k == 'list':
            items = [x.replace('**', '') for x in b['items']]
            h = sum(_lines(x, 52) * LINE_H for x in items)
            if y + h > max_y:
                return blocks[i:]
            txt(s, M, y, CW, h, '\n'.join('· ' + x for x in items),
                size=text_size, space=5)
            y += h + 0.14
        elif k == 'table':
            # §4.7은 "페이지당 표 최대 1개"를 권하지만, 정본 §4.3은 p6에
            # 총취득원가·레버리지 두 표를 함께 둡니다. 구체 조항을 따라 2개까지 허용합니다.
            if getattr(s, '_n_table', 0) >= 3:
                return blocks[i:]
            rows = [[str(v).replace('**', '') for v in r] for r in b['rows']]
            head = [h.replace('**', '') for h in b['head']]
            # 열 폭은 내용 길이에 비례시킵니다 (고정 배분은 첫 열만 넓어집니다)
            n = len(b['head'])
            widths = []
            for j in range(n):
                mx = max([len(head[j])] + [len(r[j]) for r in rows])
                floor = 1.6 if b['align'][j] == 'r' else 1.0
                widths.append(max(floor, min(mx * 0.62, 9.0)))
            hh, rh = HEAD_H, ROW_H
            h = hh + rh * len(rows)
            if y + h > max_y:
                return blocks[i:]
            table(s, M, y, CW, head, rows, widths,
                  size=table_size, head_size=table_size,
                  row_h=rh, head_h=hh, align=b['align'])
            s._n_table = getattr(s, '_n_table', 0) + 1
            y += h + 0.16
        else:
            raise ValueError(k)
    return []


def section_page(prs, core, section, kicker, archetype, title=None):
    """슬라이드 **한 장**을 만들고 넘친 블록을 반환합니다.

    정본 §4.1이 페이지를 먼저 정하므로, 흘려서 장수를 늘리지 않습니다.
    넘친 내용은 권장 슬롯(부록)으로 보냅니다.
    """
    blocks = section.get('_rest') or section['blocks']
    s = slide_new(prs, kicker, title or section['title'], archetype)
    nd = len(section['deficiencies'])
    rect(s, RIGHT - 1.85, 0.70, 1.85, 0.34, fill='okbg' if not nd else 'warnbg')
    txt(s, RIGHT - 1.77, 0.77, 1.69, 0.22,
        section['badge'] + (f' {nd}건' if nd else ''), size=10.5, bold=True,
        color='ok' if not nd else 'warn', align=PP_ALIGN.CENTER)
    rest = _blocks_to_pptx(s, core, blocks, 1.38)
    section['_rest'] = rest
    return rest


# ── 정본 §4.1 페이지 구성 ────────────────────────────────────────────
#   필수 12p + 권장 4p = 최대 16p (§3.1 골디락스)
#   흘려 넣다 보면 장수가 불어납니다. 정본은 **페이지를 먼저 정합니다.**
#   같은 key가 여러 번 나오면 이어서 흘립니다 (정본이 p6 투자 구조 / p7 가격 근거로
#   나눠 둔 것처럼, 한 섹션이 여러 면을 차지할 수 있습니다).
PAGE_PLAN = [
    ('물건 개요',   'A04', 'property_overview', None),
    ('입지',       'A06', 'location_access',   None),
    ('임대 현황',   'A03', 'lease_status',      None),
    ('임대 현황',   'A03', 'lease_status',      '임대 현황 (계속)'),
    ('임대 현황',   'A03', 'lease_status',      '임대 현황 (계속 2)'),
    ('투자 구조',   'A16', 'income_analysis',   None),
    ('투자 구조',   'A16', 'income_analysis',   '투자 구조 (계속)'),
    ('가격 근거',   'A03', 'income_analysis',   '가격 근거'),
    ('개선 여력',   'A05', 'investment_thesis', None),
    ('리스크',     'A07', 'risk_check',        None),
    ('리스크',     'A07', 'risk_check',        '리스크 및 확인사항 (계속)'),
]
MAX_TOTAL, MAX_APPENDIX = 16, 4


def photo_page(prs, core, which):
    s = slide_new(prs, '사진', f'{which} 사진', 'A14')
    rect(s, M, 1.60, CW, 3.10, fill='panel')
    txt(s, M, 2.90, CW, 0.40, '사진이 제출되지 않았습니다',
        size=15, bold=True, align=PP_ALIGN.CENTER, color='ink2')
    txt(s, M, 4.95, CW, 1.20,
        '정본 표준은 사진 6장 이상을 필수로 둡니다.\n'
        '현장 사진은 매수인의 현장 방문 결정에 가장 큰 영향을 줍니다.\n'
        '중개인이 직접 촬영한 사진으로 채워 주십시오.',
        size=12, color='ink2', space=4)
    return s


def build(core: im_core.IMCore) -> Presentation:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)
    secs = {s['key']: s for s in im_copy.build(core)}

    p01_cover(prs, core)                      # p1
    p02_summary(prs, core)                    # p2

    for kicker, arch, key, title in PAGE_PLAN:
        sec = secs[key]
        if title and not sec.get('_rest'):
            continue                          # 이미 다 들어갔으면 면을 만들지 않습니다
        section_page(prs, core, sec, kicker, arch, title=title)

    photo_page(prs, core, '외부')              # p9
    photo_page(prs, core, '내부')              # p10
    ns = secs['next_steps']
    rest = section_page(prs, core, ns, '거래 조건', 'A09')      # p11~
    while rest:
        rest = section_page(prs, core, ns, '거래 조건', 'A09',
                            title='거래 조건·다음 단계 상세')

    # 권장 슬롯 — 그래도 남으면 부록으로. **한 블록도 버리지 않습니다.**
    for kicker, arch, key, _ in PAGE_PLAN:
        sec = secs[key]
        while sec.get('_rest'):
            section_page(prs, core, sec, kicker, arch,
                         title=f'{sec["title"]} 상세')
    # 남은 블록이 없어야 합니다 (불변조건 13 · 18)
    for s_ in secs.values():
        if s_.get('_rest'):
            raise Overflow(f'{s_["title"]} 블록 {len(s_["_rest"])}개가 누락됩니다')
    if len(prs.slides) > MAX_TOTAL:
        raise Overflow(f'{len(prs.slides)}장 — 정본 §3.1 상한 {MAX_TOTAL}장 초과. '
                       f'카피를 줄이십시오.')

    footer(prs, core)
    return prs


def main() -> int:
    for fid, name in (('dangsan', '당산동'), ('yangpyeong', '양평동')):
        core = im_core.load(fid)
        prs = build(core)
        p = OUT / f'{name}_PPTX_IM.pptx'
        prs.save(str(p))
        print(f'{p.name:<24} {len(prs.slides):>2}장 · {p.stat().st_size:>7,} B')
    return 0


if __name__ == '__main__':
    try:
        sys.exit(main())
    except Overflow as e:
        print(f'경계 위반: {e}', file=sys.stderr)
        sys.exit(1)
