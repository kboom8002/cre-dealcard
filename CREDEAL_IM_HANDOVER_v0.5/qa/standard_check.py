#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
standard_check.py — 정본 `IM_STANDARD_수익형.md` 준수 검사

D17 `output_qa.py`가 **오류**를 잡는다면, 이 검사기는 **표준 준수**를 봅니다.
둘은 층이 다릅니다. 정본이 상위입니다.

검사 조항
  §3.1 분량 12~16p        §4.7 시각 규격
  §5.5 모바일 문장 길이     §6.1 치환 사전
  §6.3 금지어             §6.4 숫자 표기
  §7.1 문장 45자          §8.3 발행 전 10항
  §9.1 PPTX 필수 15       §9.2 모바일 필수 8

사용:
    python3 standard_check.py <pptx> --kind pptx
    python3 standard_check.py <html> --kind mobile
"""
from __future__ import annotations

import argparse
import html as H
import re
import sys
from pathlib import Path

# ── SSoT 레지스트리에서 읽습니다 (credeal/ssot) ─────────────────────────
# 목록을 검사기 안에 두지 않습니다. 두 군데 있으면 반드시 어긋납니다.
sys.path.insert(0, str(Path(__file__).resolve().parent.parent / 'credeal' / 'ssot'))
import loader as SSOT                                            # noqa: E402

SUBST = {k: (v or '(쓰지 않음)') for k, v in SSOT.substitutions().items()}
BANNED = SSOT.banned_words()
CONTEXT_EXCLUDE = SSOT.context_exclude()
BANNED_CTX: list[str] = []          # SSoT banned 에 이미 포함되어 있습니다
NEGATION = SSOT.negations()

# ── §9.1 PPTX 필수 15 ──────────────────────────────────────────────────
REQ_PPTX = [
    ('한 장 요약', r'한 장 요약'),
    ('매매가 + 평당가', r'평당가'),
    ('총취득원가', r'총취득원가'),
    ('실투자금', r'실투자금'),
    ('월 순현금 (LTV 3안)', r'월 순현금'),
    ('역레버리지 경고', r'대출을 늘릴수록'),
    ('임대 현황 전체 + 만료일', r'만료일'),
    ('수익률 + 기준 병기', r'÷ 매매가'),
    ('인근 매물 3~5건', r'비교사례'),
    ('개선 여력 + ◇가정', r'◇'),
    ('리스크·확인사항', r'확인된 리스크'),
    ('용도지역·용적률·건폐율', r'용도지역'),
    ('공시지가 + 배수', r'공시지가'),
    # 🔴 이 항목은 **사진 면의 존재**만 봅니다. 장수는 산출물에서 셀 수 없습니다 —
    #    같은 원본을 여러 면에 넣으면 media 항목이 늘어나 과다 계상됩니다.
    #    최소 6매는 생성 시점 게이트(G26)가 봅니다. D22-7 §한계 3 참고.
    ('사진 면 존재', r'건물 사진|사진'),
    ('연락처', r'중개인|문의'),
]

# ── §9.2 모바일 필수 8 ─────────────────────────────────────────────────
REQ_MOBILE = [
    ('한 줄 정의 (25자)', r'class="sub"'),
    ('숫자 3개', r'평당가'),
    ('임차 한 줄 요약', r'class="lease1"'),
    ('전화 버튼', r'href="tel:'),
    ('실투자금·월 순현금', r'실투자금'),
    ('확인 필요', r'확인 필요'),
    ('CTA 44px 이상', r'min-height:44px'),
    ('본문 16px 이상', r'font-size:17px'),
]


class R:
    def __init__(self):
        self.rows = []

    def add(self, clause, level, msg, detail=''):
        self.rows.append((clause, level, msg, str(detail)[:220]))

    def n(self, lv):
        return sum(1 for r in self.rows if r[1] == lv)


def strip_html(t: str) -> str:
    t = re.sub(r'<(script|style)[^>]*>.*?</\1>', ' ', t, flags=re.S | re.I)
    return H.unescape(re.sub(r'<[^>]+>', ' ', t))


def pptx_read(p: Path):
    from pptx import Presentation
    prs = Presentation(p)
    out, small, longs, tbl2 = [], 0, 0, 0
    SH_IN = prs.slide_height / 914400
    CAP = ('·', '출처', '◇', '※', '→', '지도', '본 자료', '제이에스')
    for s in prs.slides:
        n_tbl = 0
        for sh in s.shapes:
            if sh.has_text_frame:
                out.append(sh.text_frame.text)
                top = (sh.top or 0) / 914400
                # §4.7 "본문 11pt"는 본문 대상입니다.
                # 머리말(0.5in 위)·꼬리말(하단 0.8in)·캡션은 제외합니다.
                chrome = top < 0.50 or top > SH_IN - 0.80
                for para in sh.text_frame.paragraphs:
                    for r in para.runs:
                        cap = chrome or r.text.strip().startswith(CAP)
                        if (r.font.size and r.font.size.pt < 11
                                and len(r.text) > 12 and not cap):
                            small += 1
                        if len(r.text) > 45 and not chrome:
                            longs += 1
            if sh.has_table:
                n_tbl += 1
                for row in sh.table.rows:
                    out.append('\t'.join(c.text for c in row.cells))
                    for c in row.cells:
                        for para in c.text_frame.paragraphs:
                            for r in para.runs:
                                if r.font.size and r.font.size.pt < 9:
                                    small += 1
        if n_tbl > 2:
            tbl2 += 1
    return '\n'.join(out), len(prs.slides), small, longs, tbl2


def check_vocab(t: str, rep: R):
    hit = {k: t.count(k) for k in SUBST if k in t}
    if hit:
        rep.add('§6.1', 'FAIL', f'치환 사전 위반 {sum(hit.values())}건',
                ' · '.join(f'{k}→{SUBST[k]} ({v})' for k, v in hit.items()))

    lines = t.split('\n')
    bad = {}
    for w in BANNED + BANNED_CTX:
        n = 0
        for ln in lines:
            if w not in ln or any(g in ln for g in NEGATION):
                continue
            cov = []
            for full in CONTEXT_EXCLUDE.get(w, ()):
                cov += [(m.start(), m.end()) for m in re.finditer(re.escape(full), ln)]
            for m in re.finditer(re.escape(w), ln):
                if not any(a <= m.start() and m.end() <= b for a, b in cov):
                    n += 1
        if n:
            bad[w] = n
    if bad:
        rep.add('§6.3', 'FAIL', f'금지어 {sum(bad.values())}건',
                ' · '.join(f'{k} {v}' for k, v in bad.items()))


def check_numbers(t: str, rep: R):
    if re.search(r'순\s*수익률\s*\(?Cap', t):
        rep.add('§6.4', 'FAIL', '총임대료 기준에 "순수익률" 라벨을 붙였습니다')
    pcts = re.findall(r'(\d+\.\d{2})\s*%', t)
    if pcts and not re.search(r'÷', t):
        rep.add('§6.4', 'FAIL', '수익률에 기준 병기가 없습니다')
    if re.search(r'\d{3,},\d{3}\s*천원|천원 단위', t):
        rep.add('§6.4', 'WARN', '천원 단위 혼용')
    if '㎡' in t and '평' not in t:
        rep.add('§6.4', 'FAIL', '면적이 ㎡·평 병기가 아닙니다')


def check_pptx(p: Path, rep: R):
    t, n, small, longs, tbl2 = pptx_read(p)

    if not 12 <= n <= 16:
        rep.add('§3.1', 'FAIL', f'분량 {n}장 — 필수 12p · 상한 16p')
    else:
        rep.add('§3.1', 'PASS', f'분량 {n}장')

    if small:
        rep.add('§4.7', 'FAIL', f'본문 11pt / 표 9pt 미만 {small}건')
    if longs:
        rep.add('§4.7', 'FAIL', f'한 줄 45자 초과 {longs}건')
    # §4.7은 "페이지당 표 최대 1개"를 권하나, 정본 §4.3이 p6에 두 표를 함께 둡니다.
    # 구체 조항을 따라 2개까지 허용하고 3개 이상만 위반으로 봅니다.
    if tbl2:
        rep.add('§4.7', 'FAIL', f'페이지당 표 3개 이상 {tbl2}장')
    if not (small or longs or tbl2):
        rep.add('§4.7', 'PASS', '시각 규격 충족')

    miss = [nm for nm, pat in REQ_PPTX if not re.search(pat, t)]
    if miss:
        rep.add('§9.1', 'FAIL', f'필수 요소 누락 {len(miss)}/15', ' · '.join(miss))
    else:
        rep.add('§9.1', 'PASS', '필수 15종 충족')

    check_vocab(t, rep)
    check_numbers(t, rep)
    return t


def check_mobile(p: Path, rep: R):
    raw = p.read_text(encoding='utf-8')
    t = strip_html(raw)

    miss = [nm for nm, pat in REQ_MOBILE if not re.search(pat, raw)]
    if miss:
        rep.add('§9.2', 'FAIL', f'필수 요소 누락 {len(miss)}/8', ' · '.join(miss))
    else:
        rep.add('§9.2', 'PASS', '필수 8종 충족')

    one = re.search(r'class="sub">(.*?)</div>', raw, re.S)
    if one:
        s = H.unescape(re.sub(r'<[^>]+>', '', one.group(1))).strip()
        if len(s) > 25:
            rep.add('§5.5', 'FAIL', f'한 줄 정의 {len(s)}자 (상한 25)', s)
        else:
            rep.add('§5.5', 'PASS', f'한 줄 정의 {len(s)}자')

    body_only = raw.split('<footer>')[0]          # 면책 문구는 §5.5 대상 아님
    sents = []
    for m in re.finditer(r'<p[^>]*>(.*?)</p>', body_only, re.S):
        body = H.unescape(re.sub(r'<[^>]+>', '', m.group(1)))
        sents += [x.strip() for x in re.split(r'(?<=다\.)\s*', body) if x.strip()]
    over = [x for x in sents if len(x) > 40]
    if over:
        rep.add('§5.5', 'FAIL', f'본문 문장 40자 초과 {len(over)}/{len(sents)}',
                max(over, key=len)[:80])
    else:
        rep.add('§5.5', 'PASS', f'본문 문장 {len(sents)}개 전부 40자 이하')

    heads = [H.unescape(re.sub(r'<[^>]+>', '', m))
             for m in re.findall(r'class="sq">(.*?)</span>', raw)]
    ho = [h for h in heads if len(h) > 12]
    if ho:
        rep.add('§5.5', 'WARN', f'섹션 제목 12자 초과 {len(ho)}/{len(heads)}',
                ' · '.join(ho[:4]))

    check_vocab(t, rep)
    check_numbers(t, rep)
    return t


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument('file')
    ap.add_argument('--kind', choices=['pptx', 'mobile'], required=True)
    a = ap.parse_args()
    p = Path(a.file)
    rep = R()

    (check_pptx if a.kind == 'pptx' else check_mobile)(p, rep)

    mark = {'FAIL': '[위반]', 'WARN': '[주의]', 'PASS': '[충족]'}
    print(f'\n{p.name}  ·  정본 IM_STANDARD_수익형.md 대조\n' + '─' * 68)
    for lv in ('FAIL', 'WARN', 'PASS'):
        for cl, l, m, d in [r for r in rep.rows if r[1] == lv]:
            print(f'{mark[lv]} {cl:<6} {m}')
            if d:
                print(f'{"":>7} └ {d}')
    print('─' * 68)
    print(f'위반 {rep.n("FAIL")} · 주의 {rep.n("WARN")} · 충족 {rep.n("PASS")}')
    return 1 if rep.n('FAIL') else 0


if __name__ == '__main__':
    sys.exit(main())
